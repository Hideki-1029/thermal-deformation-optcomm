from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "papers/seminar/ISSLスライドテンプレート.pptx"
OUTPUT = ROOT / "papers/seminar/20260721_optcommrg_takamoto_issl_preview5.pptx"

NAVY = RGBColor(8, 43, 82)
BLUE = RGBColor(27, 116, 170)
SKY = RGBColor(91, 181, 216)
PALE = RGBColor(235, 246, 251)
PALE_GRAY = RGBColor(244, 246, 248)
GRAY = RGBColor(102, 111, 121)
INK = RGBColor(35, 43, 52)
RED = RGBColor(208, 59, 70)
ORANGE = RGBColor(236, 148, 39)
GREEN = RGBColor(35, 149, 113)
WHITE = RGBColor(255, 255, 255)
GRID = RGBColor(184, 213, 226)
FONT = "Yu Gothic"


def I(v: float):
    return Inches(v)


def remove_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def set_text(shape, text: str, size: float | None = None, color: RGBColor | None = None, bold: bool | None = None) -> None:
    shape.text = text
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.name = FONT
            if size is not None:
                run.font.size = Pt(size)
            if color is not None:
                run.font.color.rgb = color
            if bold is not None:
                run.font.bold = bold


def add_text(slide, text, x, y, w, h, *, size=28, color=INK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    sh = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = sh.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = I(0.04)
    tf.margin_top = tf.margin_bottom = I(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return sh


def add_box(slide, x, y, w, h, *, fill=WHITE, line=GRID, rounded=False, width=1.3):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, I(x), I(y), I(w), I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(width)
    return sh


def add_arrow(slide, x1, y1, x2, y2, *, color=BLUE, width=3):
    sh = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    sh.line.color.rgb = color
    sh.line.width = Pt(width)
    sh.line.end_arrowhead = True
    return sh


def add_bullets(slide, items, x, y, w, h, *, size=27, accent=BLUE):
    sh = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = sh.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = I(0.03)
    first = True
    for text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(8)
        r = p.add_run(); r.text = "• "; r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = accent; r.font.bold = True
        r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = INK
    return sh


def add_notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear(); tf.text = text


def set_category_and_page(slide, category: str, page: int) -> None:
    for sh in slide.shapes:
        if not hasattr(sh, "text_frame"):
            continue
        top = sh.top / 914400
        if top < 1.0 and sh.width / 914400 > 8.0:
            set_text(sh, category, 20, GRAY, False)
        elif top > 13.0:
            set_text(sh, str(page), 18, GRAY, False)


def clear_content(slide, keep_indices: set[int]) -> None:
    for i, shape in reversed(list(enumerate(slide.shapes))):
        if i not in keep_indices:
            remove_shape(shape)


def style_title_slide(slide) -> None:
    set_text(slide.shapes[0], "Feedforward and Adaptive Correction of\nTime-Varying Thermal Bias for Coarse Acquisition\nin Optical Communication Systems", 31, NAVY, True)
    set_text(slide.shapes[1], "Hideki Takamoto (The University of Tokyo)", 23, INK, False)
    set_text(slide.shapes[2], "2026/07/21", 22, GRAY, False)
    set_text(slide.shapes[3], "", 20, GRAY, False)
    set_text(slide.shapes[4], "Optical Communication Research Group", 22, GRAY, False)
    add_text(slide, "熱ひずみ予測と適応補正による\n光通信粗捕捉性能向上の検討", 10.58, 12.0, 11.5, 1.35, size=24, color=BLUE, bold=True)
    add_notes(slide, "7月6日の発表以降の進捗として、階層ΔTモデル、PAT評価、軌道予測誤差の導入を報告する。")


def style_section_slide(slide) -> None:
    set_text(slide.shapes[0], "01", 28, WHITE, False)
    set_text(slide.shapes[1], "Overview", 32, WHITE, True)
    add_notes(slide, "冒頭で前回までの位置づけと今回のゴールを共有する。")


def style_common_header(slide, title: str, section: str) -> None:
    set_text(slide.shapes[0], title, 32, NAVY, True)
    # category placeholder index after content deletion is 1; call before deletion.
    for sh in slide.shapes:
        if hasattr(sh, "text") and sh.text in {"スライドの基本", "Overview", "Background", "Progress since 7/6"}:
            set_text(sh, section, 20, GRAY, False)


def style_goal_slide(slide) -> None:
    # Keep title, section label, slide number; remove template body.
    clear_content(slide, {0, 2, 3})
    set_text(slide.shapes[0], "今日はモデルを見せ、仮定をレビューしてもらう", 32, NAVY, True)
    set_category_and_page(slide, "Overview", 3)

    columns = [
        ("01", "Hierarchical ΔT", "共有感度 a と\nケースDC b_case", BLUE),
        ("02", "PAT + orbit error", "非熱誤差込みの\n捕捉性能", GREEN),
        ("03", "Discussion", "Adaptiveで何を\n更新するか", ORANGE),
    ]
    for i, (num, head, body, color) in enumerate(columns):
        x = 1.25 + i * 8.25
        add_text(slide, num, x, 3.35, 1.2, 0.55, size=21, color=color, bold=True)
        add_box(slide, x, 4.05, 7.35, 4.35, fill=WHITE, line=color, width=2.0)
        band = add_box(slide, x, 4.05, 7.35, 0.68, fill=color, line=color)
        add_text(slide, head, x + 0.25, 5.05, 6.85, 0.75, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.35, 6.15, 6.65, 1.35, size=27, color=INK, align=PP_ALIGN.CENTER)

    add_text(slide, "今日特に議論したいこと", 1.25, 9.35, 6.0, 0.55, size=25, color=NAVY, bold=True)
    add_bullets(slide, [
        "aを事前同定し、b_caseを軌道上更新する分担",
        "熱と軌道誤差が同帯域にある場合の可観測性",
        "現在のscan model・誤差条件の妥当性",
    ], 1.25, 10.05, 23.8, 2.65, size=26)
    add_notes(slide, "本日は完成報告ではなく設計レビュー。モデル、PAT、Adaptiveの3点についてフィードバックをもらいたい。")


def style_background_slide(slide) -> None:
    clear_content(slide, {0, 2, 3})
    set_text(slide.shapes[0], "熱ひずみは粗捕捉開始時の scan-center error として残る", 32, NAVY, True)
    set_category_and_page(slide, "Background", 4)

    steps = [
        ("光フィードバック前", "姿勢・軌道・熱LOSを含む\n不確定領域を探索", BLUE),
        ("日照・蝕／機器発熱", "温度場 → 差動膨張 →\nSTT–LCT相対回転", ORANGE),
        ("予測FF補正", "熱LOSをscan centerから引き\n探索域・捕捉時間を低減", GREEN),
    ]
    for i, (head, body, color) in enumerate(steps):
        x = 1.2 + i * 8.35
        add_box(slide, x, 3.7, 7.1, 5.45, fill=PALE_GRAY, line=color, rounded=True, width=2.0)
        add_text(slide, head, x + 0.3, 4.25, 6.5, 0.65, size=27, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.4, 5.55, 6.3, 1.65, size=26, color=INK, align=PP_ALIGN.CENTER)
        if i == 1:
            add_text(slide, "10²–10³ µrad", x + 1.0, 7.75, 5.1, 0.65, size=29, color=RED, bold=True, align=PP_ALIGN.CENTER)
        if i < 2:
            add_arrow(slide, x + 7.35, 6.4, x + 8.05, 6.4, color=SKY, width=3)
    add_text(slide, "構造を変更せず、予測可能な熱成分だけを運用側で除去する", 3.4, 10.45, 19.8, 0.85, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "すべての指向誤差を消す手法ではない", 6.0, 11.75, 14.6, 0.6, size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "光RG向けの復習。熱変形自体ではなくSTT基準とLCT光軸の差が粗捕捉開始時のscan-center errorになる。")


def style_progress_slide(slide) -> None:
    clear_content(slide, {0, 2, 3})
    set_text(slide.shapes[0], "2週間で理想補正からケース横断モデル＋実データ外乱へ進んだ", 32, NAVY, True)
    set_category_and_page(slide, "Progress since 7/6", 5)

    add_text(slide, "7/6", 1.3, 3.35, 11.2, 0.65, size=28, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "7/21", 14.1, 3.35, 11.2, 0.65, size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_box(slide, 1.3, 4.2, 11.2, 6.8, fill=PALE_GRAY, line=GRID)
    add_box(slide, 14.1, 4.2, 11.2, 6.8, fill=PALE, line=SKY)
    left = ["熱真値補正が中心", "数ケースの予備評価", "Gaussian軌道誤差", "任意LOS横断2軸", "Adaptiveは概念"]
    right = ["階層ΔT：共有 a + b_case", "21ケース・LOO評価", "Sentinel-1 TLE vs POEORB", "STT/body座標へ射影", "b_case/DC更新へ再定義"]
    add_bullets(slide, left, 2.0, 4.85, 9.9, 5.25, size=26, accent=GRAY)
    add_bullets(slide, right, 14.8, 4.85, 9.9, 5.25, size=26, accent=BLUE)
    add_arrow(slide, 12.75, 7.5, 13.65, 7.5, color=SKY, width=4)
    add_text(slide, "現在の問い：同帯域の非熱誤差と共存しても、軽量熱モデルで捕捉性能を改善できるか", 2.1, 11.75, 22.5, 0.85, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "7月6日との違いを一枚で示す。当時の25.8秒は条件が異なる予備結果なので今回とは直接比較しない。")


def build() -> Presentation:
    prs = Presentation(str(TEMPLATE))
    # Retain the first five template slides and their original ISSL layouts.
    for slide_id in list(prs.slides._sldIdLst)[5:]:
        prs.slides._sldIdLst.remove(slide_id)

    style_title_slide(prs.slides[0])
    style_section_slide(prs.slides[1])
    style_goal_slide(prs.slides[2])
    style_background_slide(prs.slides[3])
    style_progress_slide(prs.slides[4])
    return prs


def main() -> None:
    prs = build()
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
