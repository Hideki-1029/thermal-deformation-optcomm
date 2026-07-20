from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_optcommrg_issl_preview5 import (
    TEMPLATE,
    NAVY,
    BLUE,
    SKY,
    PALE,
    PALE_GRAY,
    GRAY,
    INK,
    RED,
    ORANGE,
    GREEN,
    WHITE,
    GRID,
    FONT,
    I,
    remove_shape,
    set_text,
    add_text,
    add_box,
    add_arrow,
    add_bullets,
    add_notes,
    style_title_slide,
    style_section_slide,
    style_goal_slide,
    style_background_slide,
    style_progress_slide,
)


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "papers/seminar/20260721_optcommrg_takamoto_issl.pptx"
PURPLE = RGBColor(93, 64, 132)
PALE_RED = RGBColor(253, 239, 240)
PALE_ORANGE = RGBColor(254, 245, 230)
PALE_GREEN = RGBColor(234, 248, 241)


def clear_all(slide) -> None:
    for sh in reversed(list(slide.shapes)):
        remove_shape(sh)


def reset_content(slide, title: str, category: str, page: str | int) -> None:
    clear_all(slide)
    add_text(slide, category, 0.60, 0.24, 17.4, 0.55, size=20, color=GRAY)
    add_text(slide, title, 1.68, 1.42, 24.1, 1.05, size=32, color=NAVY, bold=True)
    add_text(slide, str(page), 23.8, 14.02, 2.2, 0.42, size=18, color=GRAY, align=PP_ALIGN.RIGHT)


def reset_section(slide, number: str, title: str, page: int) -> None:
    clear_all(slide)
    add_text(slide, number, 10.6, 4.50, 5.4, 0.65, size=27, color=WHITE, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, I(11.45), I(5.72), I(3.7), I(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = WHITE; line.line.fill.background()
    add_text(slide, title, 2.0, 6.35, 22.7, 1.15, size=33, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, str(page), 23.8, 14.02, 2.2, 0.42, size=18, color=WHITE, align=PP_ALIGN.RIGHT)
    add_notes(slide, f"Section {number}: {title}")


def new_content(prs: Presentation, title: str, category: str, page: str | int | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    reset_content(slide, title, category, page if page is not None else len(prs.slides))
    return slide


def new_section(prs: Presentation, number: str, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    reset_section(slide, number, title, len(prs.slides))
    return slide


def add_table(slide, headers, rows, x, y, w, h, *, widths=None, font_size=22, accent_col=None, header=NAVY):
    sh = slide.shapes.add_table(len(rows) + 1, len(headers), I(x), I(y), I(w), I(h))
    table = sh.table
    if widths:
        for col, cw in zip(table.columns, widths):
            col.width = I(cw)
    for c, text in enumerate(headers):
        cell = table.cell(0, c); cell.text = str(text)
        cell.fill.solid(); cell.fill.fore_color.rgb = header
    for r, row in enumerate(rows, 1):
        for c, text in enumerate(row):
            cell = table.cell(r, c); cell.text = str(text)
            cell.fill.solid(); cell.fill.fore_color.rgb = PALE if c == accent_col else (WHITE if r % 2 else PALE_GRAY)
    for r in range(len(rows) + 1):
        for c in range(len(headers)):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = I(0.07)
            cell.margin_top = cell.margin_bottom = I(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT; run.font.size = Pt(font_size)
                    run.font.bold = r == 0; run.font.color.rgb = WHITE if r == 0 else INK
    return sh


def add_picture_fit(slide, path: Path, x, y, w, h):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, I(x), I(y), I(w), I(h))
    frame.fill.background(); frame.line.color.rgb = GRID; frame.line.width = Pt(1)
    if not path.exists():
        add_text(slide, f"Missing: {path.name}", x + 0.3, y + 0.3, w - 0.6, 0.8, size=22, color=RED)
        return
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    slide.shapes.add_picture(str(path), I(x + (w - pw) / 2), I(y + (h - ph) / 2), I(pw), I(ph))


def add_process(slide, labels, x, y, total_w, *, colors=None, box_h=1.35):
    colors = colors or [NAVY, BLUE, GREEN, ORANGE, PURPLE]
    gap = 0.42
    bw = (total_w - gap * (len(labels) - 1)) / len(labels)
    for i, label in enumerate(labels):
        xx = x + i * (bw + gap)
        add_box(slide, xx, y, bw, box_h, fill=colors[i % len(colors)], line=colors[i % len(colors)])
        add_text(slide, label, xx + 0.08, y + 0.12, bw - 0.16, box_h - 0.24, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(labels) - 1:
            add_arrow(slide, xx + bw + 0.06, y + box_h / 2, xx + bw + gap - 0.06, y + box_h / 2, color=SKY, width=2)


def add_metric(slide, value, label, x, y, w, *, color=BLUE):
    add_box(slide, x, y, w, 1.55, fill=WHITE, line=color, rounded=True, width=1.8)
    add_text(slide, value, x + 0.1, y + 0.14, w - 0.2, 0.62, size=34, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, y + 0.88, w - 0.2, 0.4, size=19, color=GRAY, align=PP_ALIGN.CENTER)


def add_bars(slide, labels, values, x, y, w, h, *, colors, suffix="", max_value=None):
    max_value = max_value or max(values)
    row_h = h / len(values); label_w = w * 0.31; bar_w = w - label_w - 1.15
    for i, (label, value, color) in enumerate(zip(labels, values, colors)):
        yy = y + i * row_h
        add_text(slide, label, x, yy + 0.04, label_w - 0.1, row_h - 0.08, size=21, valign=MSO_ANCHOR.MIDDLE)
        add_box(slide, x + label_w, yy + 0.15, bar_w, row_h - 0.30, fill=PALE_GRAY, line=PALE_GRAY)
        add_box(slide, x + label_w, yy + 0.15, max(0.08, bar_w * value / max_value), row_h - 0.30, fill=color, line=color)
        add_text(slide, f"{value:g}{suffix}", x + label_w + bar_w + 0.12, yy, 1.0, row_h, size=20, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)


def slide_architecture(slide, page: int) -> None:
    reset_content(slide, "熱構造解析から粗捕捉補正までを一気通貫で評価する", "Overview", page)
    add_process(slide, ["Orbit / attitude\nheat & coating", "Thermal Desktop\ntemperature", "Femap / Nastran\nrotation", "STT–LCT\nthermal LOS", "Lightweight\nmodel", "PAT\nscan center"], 1.0, 3.35, 24.7)
    add_box(slide, 1.35, 5.5, 11.4, 5.25, fill=PALE_GRAY, line=GRID)
    add_text(slide, "High-fidelity reference", 1.8, 5.95, 10.5, 0.6, size=27, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["軌道熱環境・日照／蝕", "温度場・熱変形", "STT–LCT相対回転"], 2.0, 7.0, 9.9, 2.9, size=26)
    add_box(slide, 13.8, 5.5, 11.4, 5.25, fill=PALE, line=SKY)
    add_text(slide, "Operational layer", 14.25, 5.95, 10.5, 0.6, size=27, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["温度センサ＋運用フラグ", "予測LOSをscan centerへ", "捕捉後残差でモデル更新"], 14.45, 7.0, 9.9, 2.9, size=26)
    add_text(slide, "TD/Femapはオンボード実装ではなく、軽量モデルの教師・評価用真値", 4.0, 11.7, 18.7, 0.7, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "全体フロー。TD/Femapは高忠実度参照で、軌道上では軽量モデルのみを使う。")


def slide_claim(slide, page: int) -> None:
    reset_content(slide, "熱の予測精度ではなく、粗捕捉性能の改善までを主張する", "Overview", page)
    add_box(slide, 1.3, 3.35, 7.3, 6.8, fill=PALE_GRAY, line=GRID)
    add_text(slide, "1", 1.75, 3.75, 1.0, 0.6, size=25, color=BLUE, bold=True)
    add_text(slide, "Quantify", 2.1, 4.55, 5.7, 0.7, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "熱構造解析による\nSTT–LCT相対LOS", 2.0, 6.1, 5.9, 1.4, size=27, align=PP_ALIGN.CENTER)
    add_box(slide, 9.65, 3.35, 7.3, 6.8, fill=PALE, line=SKY)
    add_text(slide, "2", 10.1, 3.75, 1.0, 0.6, size=25, color=BLUE, bold=True)
    add_text(slide, "Predict", 10.45, 4.55, 5.7, 0.7, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "共有感度＋ケースDCの\n階層ΔTモデル", 10.35, 6.1, 5.9, 1.4, size=27, align=PP_ALIGN.CENTER)
    add_box(slide, 18.0, 3.35, 7.3, 6.8, fill=PALE_GREEN, line=GREEN)
    add_text(slide, "3", 18.45, 3.75, 1.0, 0.6, size=25, color=GREEN, bold=True)
    add_text(slide, "Improve", 18.8, 4.55, 5.7, 0.7, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "scan center補正による\n成功率・捕捉時間改善", 18.7, 6.1, 5.9, 1.4, size=27, align=PP_ALIGN.CENTER)
    add_text(slide, "すべての指向誤差ではなく、予測可能な熱成分を除去する", 4.3, 11.4, 18.0, 0.8, size=29, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, "研究の主張を3段階で示す。最終評価はLOS RMSEだけでなく捕捉時間。")


def build() -> Presentation:
    prs = Presentation(str(TEMPLATE))
    # Repurpose all 18 template slides before appending new slides. This keeps
    # the original ISSL masters/layouts and avoids duplicate slide-part names.

    style_title_slide(prs.slides[0])
    style_section_slide(prs.slides[1])
    style_goal_slide(prs.slides[2])
    style_background_slide(prs.slides[3])
    style_progress_slide(prs.slides[4])
    slide_architecture(prs.slides[5], 6)
    slide_claim(prs.slides[6], 7)
    reset_section(prs.slides[7], "02", "Thermo-structural analysis", 8)

    # 9 spacecraft
    s = prs.slides[8]; reset_content(s, "箱型衛星モデルで STT–LCT間の熱変形を評価する", "Thermo-structural analysis", 9)
    add_box(s, 1.3, 3.25, 11.0, 7.6, fill=PALE_GRAY, line=GRID)
    add_box(s, 4.0, 5.0, 5.7, 3.7, fill=PALE, line=BLUE)
    add_text(s, "PZ: STT", 5.2, 4.15, 3.2, 0.55, size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "MZ: LCT\nboresight ≈ −Z", 4.8, 6.0, 4.1, 1.1, size=26, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "MY: PCDU", 1.8, 6.0, 2.2, 0.5, size=22, color=ORANGE, bold=True)
    add_text(s, "PY: PROP", 9.75, 6.0, 2.2, 0.5, size=22, color=ORANGE, bold=True)
    add_arrow(s, 6.85, 8.0, 6.85, 9.6, color=GREEN, width=4)
    add_table(s, ["Item", "Setting"], [("Bus size", "590 × 600 × 990 mm"), ("Panels", "10 mm shell"), ("Material", "A5052"), ("LCT", "MZ center"), ("STT", "PZ center")], 13.1, 3.25, 12.1, 7.2, widths=[4.3, 7.8], font_size=23, accent_col=1)
    add_text(s, "太陽面 MX / MY / PX / PY によって支配軸と符号が変わる", 3.4, 11.65, 19.8, 0.65, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "簡略shell衛星。LCTはMZ、STTはPZ、PROP/PCDUは±Y側。")

    # 10 LOS
    s = prs.slides[9]; reset_content(s, "PAT入力には STT基準で見たLCT光軸の相対回転を用いる", "Thermo-structural analysis", 10)
    add_picture_fit(s, ROOT / "results/femap_deformation/15_LTAN06_800km_1213COLD_MY_STTLCT_HEAT_MY_0p5/los_definition_comparison.png", 1.1, 3.05, 15.2, 9.0)
    add_box(s, 17.0, 3.05, 8.2, 9.0, fill=PALE, line=SKY, rounded=True)
    add_text(s, "採用：far_field_los", 17.4, 3.55, 7.4, 0.65, size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["LCT光軸回転 − STT姿勢基準回転", "遠方通信のscan centerへ直接寄与", "代表点間並進のcenterline tiltは加えない", "stt_relative_losは診断用"], 17.5, 4.8, 7.0, 4.7, size=24)
    add_text(s, "θthermal,true(t) = [θx, θy]", 17.6, 10.35, 6.9, 0.6, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "PAT主入力はfar-field relative rotation。並進成分を含む角度は診断用。")

    # 11 orbit thermal
    s = prs.slides[10]; reset_content(s, "基本ケースは LTAN06・800 kmを3軌道、約60秒刻みで解析する", "Thermo-structural analysis", 11)
    add_table(s, ["Parameter", "Setting"], [("Orbit", "LTAN06 / 800 km / COLD"), ("Period", "≈6050 s (101 min)"), ("Duration", "≈18157 s (3 orbits)"), ("Sampling", "≈60.5 s / 301 samples"), ("Eclipse", "TD LOGIC_SUN"), ("Sun face", "MX / MY / PX / PY")], 1.25, 3.05, 11.8, 7.9, widths=[4.2, 7.6], font_size=23, accent_col=1)
    add_box(s, 13.8, 3.05, 11.4, 7.9, fill=PALE, line=SKY, rounded=True)
    add_text(s, "Additional environments", 14.25, 3.55, 10.5, 0.6, size=27, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["HOT thermal environment", "Black / Alodine / α=ε=0.5", "LTAN18 / 693 km Sentinel-1 proxy", "case25: nearly all-sun"], 14.5, 4.75, 9.9, 4.3, size=26)
    add_text(s, "蝕入り・蝕明けはTD出力LOGIC_SUNを正本とする", 4.7, 11.65, 17.2, 0.7, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "基本3軌道。先頭1軌道をモデル同定、残り2軌道をtestに用いる。")

    # 12 structural
    s = prs.slides[11]; reset_content(s, "簡略shellモデルを軽量モデル評価用の高忠実度参照とする", "Thermo-structural analysis", 12)
    add_table(s, ["Structural parameter", "Value"], [("Material", "Aluminum 5052"), ("Young's modulus", "70.327 GPa"), ("Poisson ratio", "0.33"), ("CTE", "2.376×10⁻⁵ /°C"), ("Density", "2685 kg/m³"), ("Reference temperature", "23.9°C"), ("Constraint", "small region near STT")], 1.25, 3.05, 11.9, 8.5, widths=[5.1, 6.8], font_size=22, accent_col=1)
    add_picture_fit(s, ROOT / "results/femap_deformation/15_LTAN06_800km_1213COLD_MY_STTLCT_HEAT_MY_0p5/stt_lct_motion_overview.png", 13.8, 3.05, 11.4, 8.5)
    add_text(s, "TD mapper温度 → 各時刻の熱荷重 → STT/LCT基準面の変位・回転", 3.0, 12.0, 20.6, 0.65, size=27, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "実機詳細モデルではない。構造・拘束条件依存性が残る点を明示する。")

    # 13 heat
    s = prs.slides[12]; reset_content(s, "内部発熱は面間温度差とケースDCの両方を変える", "Thermo-structural analysis", 13)
    add_table(s, ["Component", "Heat", "Location"], [("LCT", "10 W", "MZ center"), ("STT", "1.5 W", "PZ center"), ("PROP", "25 W", "PY side"), ("PCDU", "10 W", "MY side")], 1.25, 3.1, 10.4, 5.7, widths=[3.4, 2.8, 4.2], font_size=23, accent_col=1)
    add_table(s, ["Mode", "LCT/STT", "PROP", "PCDU"], [("STTLCT", "ON", "OFF", "OFF"), ("+PROP", "ON", "25 W", "OFF"), ("+PCDU", "ON", "OFF", "10 W"), ("ALL", "ON", "25 W", "10 W"), ("Half power", "ON", "12.5 W", "10 W")], 12.4, 3.1, 12.8, 6.6, widths=[4.3, 2.9, 2.8, 2.8], font_size=21, accent_col=0)
    add_box(s, 2.1, 10.45, 22.3, 1.8, fill=PALE_ORANGE, line=ORANGE, rounded=True)
    add_text(s, "発熱の主効果はΔTへ、ΔTに入りきらない残差DCは b_caseへ", 2.6, 10.95, 21.3, 0.65, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "発熱モードはLevel 2の説明変数。半電力case22はON/OFFフラグの限界を確認する。")

    # 14 cases
    s = prs.slides[13]; reset_content(s, "21ケースで太陽面・発熱・被覆・軌道を横断評価する", "Thermo-structural analysis", 14)
    add_table(s, ["Varied factor", "Cases", "Purpose"], [("Sun face", "MX / MY / PX / PY", "axis & sign"), ("Power mode", "STTLCT / +PROP / +PCDU / ALL", "ΔT & DC"), ("Coating", "0.5 / Black / Alodine", "amplitude & floor"), ("Environment", "COLD / HOT", "case DC"), ("Power level", "PROP 25 / 12.5 W", "ON/OFF limit"), ("Orbit", "LTAN06 / LTAN18", "thermal history")], 1.2, 3.05, 24.0, 7.7, widths=[4.5, 10.1, 9.4], font_size=23, accent_col=2)
    add_box(s, 2.0, 11.25, 22.5, 1.35, fill=PALE_RED, line=RED, rounded=True)
    add_text(s, "除外：01–03, 07（MZ太陽指向）— 現在の1軸sunfaceモデル対象外", 2.5, 11.62, 21.5, 0.55, size=26, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "21ケースはいずれも同じ基本構造の感度ケースで、別機体への汎化ではない。")

    # 15 data flow
    s = prs.slides[14]; reset_content(s, "case_idを軸に解析条件と成果物を一貫管理する", "Thermo-structural analysis", 15)
    add_process(s, ["case_matrix.xlsx\norbit_catalog", "TD temperature\nLOGIC_SUN", "mapper\noutput.dat", "Femap\nrotation", "los_angles.csv", "lightweight\ndataset"], 1.0, 3.2, 24.7)
    add_box(s, 1.3, 5.55, 24.0, 5.25, fill=PALE_GRAY, line=GRID, rounded=True)
    add_text(s, "自動化・再現性", 1.8, 6.0, 5.0, 0.6, size=27, color=NAVY, bold=True)
    add_bullets(s, ["case_idから下流パスを規則解決", "温度場チェック後にFemap・LOSへ", "同じIDで温度・LOS・PATを追跡", "cases 4–25を同一手順で再評価"], 1.9, 7.0, 22.5, 3.0, size=26)
    add_text(s, "ケース数より、同じ条件を再現して比較できることを優先", 4.4, 11.65, 17.8, 0.7, size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "光学特性未反映やmapper混線を経験し、ケースID管理を強化した。")

    # 16 representative
    s = prs.slides[15]; reset_content(s, "代表ケースでは面間温度差と熱LOSが同じ軌道周期で変動する", "Thermo-structural results", 16)
    add_picture_fit(s, ROOT / "results/femap_deformation/04_LTAN06_800km_1213COLD_MY_ALL_HEAT_MY_0p5/default_surface_9points_temperature_overview.png", 1.1, 3.0, 11.8, 8.6)
    add_picture_fit(s, ROOT / "results/femap_deformation/04_LTAN06_800km_1213COLD_MY_ALL_HEAT_MY_0p5_far_field_los_angle_budget.png", 13.6, 3.0, 11.8, 8.6)
    add_text(s, "Temperature", 4.1, 11.85, 5.8, 0.55, size=23, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "Far-field LOS", 16.6, 11.85, 5.8, 0.55, size=23, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "case04。太陽面と反対面の温度差が支配軸LOSを説明できそうだと分かった。")

    # 17 sensitivity
    s = prs.slides[16]; reset_content(s, "太陽面で支配軸と符号が決まり、被覆・発熱は振幅とDCを変える", "Thermo-structural results", 17)
    for i, (face, axis, sign, val, color) in enumerate([("MX", "X", "+", 30.6, GREEN), ("PX", "X", "−", 28.1, BLUE), ("MY", "Y", "+", 28.6, ORANGE), ("PY", "Y", "−", 28.7, RED)]):
        x = 1.2 + i * 6.15
        add_box(s, x, 3.15, 5.45, 4.9, fill=WHITE, line=color, rounded=True, width=2)
        add_text(s, face, x + 0.2, 3.55, 5.05, 0.65, size=32, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, f"dominant {axis}-axis", x + 0.3, 4.6, 4.85, 0.5, size=22, color=GRAY, align=PP_ALIGN.CENTER)
        add_text(s, sign, x + 1.5, 5.35, 2.45, 0.8, size=43, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, f"|a|≈{val:.1f} µrad/°C", x + 0.25, 6.65, 4.95, 0.55, size=23, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_table(s, ["Factor", "Observed effect"], [("Black coating", "time-varying floor ↑"), ("PROP / PCDU", "ΔT + case DC"), ("HOT / all-sun", "nearly static LOS"), ("Half power", "intermediate b")], 3.2, 8.75, 20.2, 3.3, widths=[7.0, 13.2], font_size=21, accent_col=1)
    add_notes(s, "対向面で符号反転。被覆は残差床、発熱はΔTとDC、軌道条件は時変量を変える。")

    # Section 03: repurpose the template's final blue slide.
    reset_section(prs.slides[17], "03", "Lightweight LOS model", 18)

    s = new_content(prs, "軌道上モデルには少数入力・解釈可能性・PAT直結性が必要", "Lightweight LOS model")
    add_box(s, 1.2, 3.1, 8.0, 7.8, fill=PALE_GRAY, line=GRID, rounded=True)
    add_text(s, "Operational requirements", 1.6, 3.55, 7.2, 0.6, size=27, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["TD/Femapをオンボードで回さない", "少数温度センサ", "運用モード変更に対応", "固定少数係数", "scan-center補正を直接出力"], 1.8, 4.65, 6.8, 5.2, size=25)
    add_table(s, ["Candidate", "Pros / cons"], [("Static bias", "lightweight / no dynamics"), ("Fourier", "periodic / orbit dependent"), ("Generic temperature", "flexible / collinear"), ("Sunface ΔT", "physical / few inputs"), ("Hierarchical ΔT", "dynamics + case DC")], 10.0, 3.1, 15.2, 7.8, widths=[5.7, 9.5], font_size=23, accent_col=0)
    add_text(s, "本命：太陽面−反対面温度差 ΔT と運用フラグ", 4.2, 11.6, 18.2, 0.7, size=29, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Fourierよりも物理入力で軌道条件の変化を扱いたい。")

    s = new_content(prs, "コンポ温度の時系列追加は共線と低SNRで係数が壊れる", "Lightweight LOS model")
    add_box(s, 1.3, 3.1, 11.3, 7.8, fill=PALE_RED, line=RED, rounded=True)
    add_text(s, "Attempt", 1.8, 3.6, 10.3, 0.6, size=28, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "LOS = b + a·ΔT\n+ c·(Tattach−Tref)", 2.0, 4.85, 9.9, 1.3, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["corr(ΔT,TPCDU)≈0.995", "local std≈0.1°C", "係数が巨大化", "ケース間で不安定"], 2.0, 6.8, 9.9, 3.0, size=25, accent=RED)
    add_arrow(s, 13.0, 6.8, 14.4, 6.8, color=SKY, width=4)
    add_box(s, 14.8, 3.1, 10.4, 7.8, fill=PALE_GREEN, line=GREEN, rounded=True)
    add_text(s, "Interpretation", 15.3, 3.6, 9.4, 0.6, size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "効いていたのは\n軌道内変動ではなく\nケース平均DC", 15.7, 5.0, 8.6, 2.2, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "→ b_case側へ分離", 16.1, 8.25, 7.8, 0.65, size=29, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "RMSEだけでなく、係数の安定性と運用解釈でモデルを選ぶ", 4.2, 11.65, 18.2, 0.7, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "見かけのRMSE改善より、共線で係数解釈が壊れることを問題視した。")

    s = new_content(prs, "Level 1では軌道内時変を面間温度差一本で表す", "Hierarchical ΔT model")
    add_box(s, 1.4, 3.15, 23.8, 2.7, fill=PALE, line=SKY, rounded=True)
    add_text(s, "θdom(t) ≈ bcase + asunface · ΔT(t)", 2.0, 3.95, 22.6, 0.9, size=39, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "ΔT(t) = Tsunface(t) − Topposite(t)", 5.2, 6.45, 16.3, 0.65, size=29, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    for i, (head, body, color) in enumerate([("a(sun)", "差動膨張・曲げ感度\n面ごとに共有", BLUE), ("b_case", "ΔTに入りきらない\nケース内定数", ORANGE), ("dominant axis", "MX/PX→X\nMY/PY→Y", GREEN)]):
        x = 1.6 + i * 8.1
        add_box(s, x, 8.0, 7.2, 3.5, fill=WHITE, line=color, rounded=True, width=1.8)
        add_text(s, head, x + 0.3, 8.45, 6.6, 0.6, size=28, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, body, x + 0.4, 9.45, 6.4, 1.3, size=25, align=PP_ALIGN.CENTER)
    add_notes(s, "Level 1は支配軸のみ。非支配軸は現状モデル外。")

    s = new_content(prs, "Level 2ではケースDCを太陽面と発熱モードで説明する", "Hierarchical ΔT model")
    add_box(s, 1.4, 3.15, 23.8, 2.7, fill=PALE_ORANGE, line=ORANGE, rounded=True)
    add_text(s, "bcase ≈ b0(sunface) + cPROP IPROP + cPCDU IPCDU", 1.9, 3.95, 22.8, 0.9, size=36, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_table(s, ["Case-dependent input", "Cross-case fixed parameter"], [("ΔT(t) timeseries", "a(MX/MY/PX/PY) ×4"), ("sun_face", "b0(MX/MY/PX/PY) ×4"), ("PROP/PCDU ON/OFF", "cPROP, cPCDU ×2")], 2.0, 6.8, 22.6, 4.2, widths=[11.3, 11.3], font_size=24, accent_col=1)
    add_text(s, "合計10個のスカラー係数で21ケースを表現", 5.2, 11.65, 16.3, 0.7, size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "現在は発熱フラグをMY/PYで有効化し、MX/PXでは差が小さいため無効。")

    s = new_content(prs, "先頭1軌道で時変、ケース横断LOOでDCを評価する", "Hierarchical ΔT model")
    add_process(s, ["Per case\nfirst orbit", "Fit aemp,bemp\nRidge≈OLS", "Median aemp\nper face", "Level-2 fit\nacross cases", "LOO bpred\nnext 2 orbits"], 1.1, 3.2, 24.5)
    add_table(s, ["Layer", "Train", "Test / validation"], [("Level 1", "first orbit / case", "next 2 orbits"), ("Level 2 b", "other 20 cases", "leave-one-case-out"), ("Shared a", "median of all aemp", "not full LOO")], 1.8, 6.1, 23.0, 4.2, widths=[4.4, 8.8, 9.8], font_size=23, accent_col=2)
    add_box(s, 2.2, 11.1, 22.2, 1.35, fill=PALE_RED, line=RED, rounded=True)
    add_text(s, "同一基本構造の感度ケース：別機体への汎化を示す評価ではない", 2.7, 11.48, 21.2, 0.55, size=26, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "bはLOOだがa_sharedは対象ケースを含む中央値で完全LOOではない。")

    s = new_content(prs, "共有感度 a は21ケース追加後も ±28–31 µrad/°Cで安定する", "Hierarchical ΔT model")
    add_picture_fit(s, ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/bcase_a_emp_by_sunface.png", 1.1, 3.0, 16.2, 8.8)
    add_table(s, ["Sun", "a_shared"], [("MX", "+30.6"), ("MY", "+28.6"), ("PX", "−28.1"), ("PY", "−28.7")], 18.0, 3.2, 7.2, 5.2, widths=[3.1, 4.1], font_size=24, accent_col=1)
    add_box(s, 18.0, 9.0, 7.2, 2.4, fill=PALE_GREEN, line=GREEN, rounded=True)
    add_text(s, "事前解析・熱真空試験で\n固定できる可能性", 18.4, 9.45, 6.4, 1.3, size=26, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "aの安定性がモデルの核心。Adaptiveで毎回aを学ぶ必要性は低そう。")

    s = new_content(prs, "ケースDC b は太陽面＋発熱フラグで LOO RMSE 3.8 µradとなる", "Hierarchical ΔT model")
    add_picture_fit(s, ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/bcase_b_emp_vs_b_pred.png", 1.0, 3.0, 16.0, 8.7)
    add_table(s, ["Feature", "µrad"], [("b0_MX", "+15.7"), ("b0_MY", "+2.8"), ("b0_PX", "−12.0"), ("b0_PY", "−24.0"), ("c_PROP", "−22.9"), ("c_PCDU", "−10.2")], 17.6, 3.15, 7.6, 6.7, widths=[4.1, 3.5], font_size=22, accent_col=1)
    add_metric(s, "3.80 µrad", "LOO b RMSE", 17.7, 10.25, 7.4, color=GREEN)
    add_notes(s, "最大LOO残差は半電力case22で13.2 µrad。")

    s = new_content(prs, "標準ケースは数µradまで予測し、半電力ケースでDCずれが残る", "Hierarchical ΔT model")
    add_picture_fit(s, ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/timeseries/case08_bcase_true_vs_pred.png", 1.0, 3.0, 12.1, 8.5)
    add_picture_fit(s, ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/timeseries/case22_bcase_true_vs_pred.png", 13.55, 3.0, 12.1, 8.5)
    add_text(s, "Case 08: standard / large swing", 2.8, 11.75, 8.6, 0.5, size=22, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "Case 22: PROP 12.5 W", 15.8, 11.75, 7.5, 0.5, size=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "半電力ではaは合うが、ON/OFF Level 2では中間bを表せない。")

    s = new_content(prs, "生LOS数百µradに対し、LOO test残差は平均約5.5 µradとなる", "Hierarchical ΔT model")
    add_picture_fit(s, ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/bcase_raw_vs_model_rmse.png", 1.0, 3.0, 18.3, 8.8)
    add_metric(s, "615 µrad", "median raw RMS", 19.8, 3.4, 5.3, color=BLUE)
    add_metric(s, "5.5 µrad", "mean LOO test RMSE", 19.8, 5.8, 5.3, color=GREEN)
    add_metric(s, "1–2 orders", "thermal reduction", 19.8, 8.2, 5.3, color=ORANGE)
    add_text(s, "raw RMSとprediction RMSEは異なる統計量：オーダー比較として表示", 4.0, 12.15, 18.7, 0.5, size=23, color=GRAY, align=PP_ALIGN.CENTER)
    add_notes(s, "モデル精度の主数字。統計量の違いを明示する。")

    s = new_content(prs, "新規性は式形ではなく、衛星バス相対LOSと粗捕捉への接続にある", "Positioning & limitations")
    add_box(s, 1.2, 3.1, 11.7, 8.2, fill=PALE_GRAY, line=GRID, rounded=True)
    add_text(s, "JANUS / prior art", 1.7, 3.55, 10.7, 0.6, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["Optical head内で LOS≈KΔT", "温度センサによる一次補正", "比例関係自体は既知"], 1.9, 4.8, 10.2, 3.2, size=26)
    add_text(s, "式に切片を足すこと自体を\n新規性とは主張しない", 2.1, 8.7, 9.8, 1.2, size=26, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_box(s, 13.65, 3.1, 11.7, 8.2, fill=PALE, line=SKY, rounded=True)
    add_text(s, "This work", 14.15, 3.55, 10.7, 0.6, size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["STT–LCTの衛星バス相対LOS", "太陽面−反対面ΔTの共有感度", "発熱DCの階層モデル", "scan-center FFと捕捉時間"], 14.35, 4.8, 10.2, 4.2, size=26)
    add_text(s, "限界：MZ・非支配軸・別構造・飛行実証は未評価", 4.1, 11.95, 18.5, 0.6, size=27, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "ΔT一次関係はJANUSにある。差分はバス相対LOS、階層DC、粗捕捉性能への接続。")

    s = new_content(prs, "aは事前同定、b_caseは運用条件・軌道上残差で更新する価値が高い", "Operational interpretation")
    add_table(s, ["Parameter", "Observed behavior", "Operational treatment"], [("a(sun)", "stable across cases", "ground analysis / thermal-vac test"), ("b_case", "changes with heat mode", "mode table + online update"), ("time-varying residual", "few to 10+ µrad", "slow correction if observable"), ("nonthermal error", "same orbit-frequency family", "separate state / prior")], 1.2, 3.2, 24.0, 7.2, widths=[5.0, 8.3, 10.7], font_size=23, accent_col=2)
    add_box(s, 2.3, 11.05, 22.0, 1.4, fill=PALE_GREEN, line=GREEN, rounded=True)
    add_text(s, "Adaptiveの第一候補は a ではなく b_case / model DC", 3.0, 11.45, 20.6, 0.55, size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "モデル結果からAdaptiveの対象を再整理する。")

    # Section 04
    new_section(prs, "04", "PAT evaluation and orbit error")

    s = new_content(prs, "予測LOSをscan centerから引き、残差をスキャン対象とする", "PAT evaluation")
    add_box(s, 1.4, 3.1, 23.8, 2.35, fill=PALE, line=SKY, rounded=True)
    add_text(s, "escan(t) = enonthermal(t) + θthermal,true(t) − θ̂thermal(t)", 1.9, 3.78, 22.8, 0.75, size=36, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_process(s, ["Thermal LOS\ntruth", "Model\nprediction", "Scan-center\nresidual", "Rectangular\nspiral", "Acquisition\ntime"], 1.4, 6.0, 23.8)
    add_table(s, ["Comparison", "Purpose"], [("no correction", "thermal baseline"), ("static bias", "DC baseline"), ("bcase", "proposed model"), ("thermal truth", "ideal upper bound"), ("thermal+nonthermal, no corr.", "realism baseline"), ("bcase+nonthermal", "main result")], 3.6, 8.4, 19.4, 4.0, widths=[9.7, 9.7], font_size=20, accent_col=0)
    add_notes(s, "非熱誤差は補正対象外で残る。熱モデル能力とシステム性能を分けて読む。")

    s = new_content(prs, "粗捕捉は矩形スパイラルを離散点で抽象化する", "PAT simulation conditions")
    add_box(s, 1.2, 3.1, 10.2, 8.5, fill=PALE_GRAY, line=GRID, rounded=True)
    cx, cy, scale = 6.2, 7.1, 0.72
    pts = [(0,0),(1,0),(1,1),(-1,1),(-1,-1),(2,-1),(2,2),(-2,2),(-2,-2),(3,-2),(3,3)]
    for i in range(len(pts)-1):
        add_arrow(s, cx+pts[i][0]*scale, cy-pts[i][1]*scale, cx+pts[i+1][0]*scale, cy-pts[i+1][1]*scale, color=BLUE, width=2)
    for px, py in pts:
        d=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,I(cx+px*scale-0.08),I(cy-py*scale-0.08),I(0.16),I(0.16)); d.fill.solid(); d.fill.fore_color.rgb=BLUE; d.line.fill.background()
    t=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,I(cx+1.55),I(cy-1.75),I(0.5),I(0.5)); t.fill.solid(); t.fill.fore_color.rgb=PALE_ORANGE; t.line.color.rgb=ORANGE
    add_table(s, ["Parameter", "Value"], [("Range", "±1600 µrad"), ("Step", "40 µrad"), ("Detection radius", "25 µrad"), ("Dwell", "0.1 s / point"), ("Points", "6561"), ("Max nominal time", "656.1 s")], 12.2, 3.1, 13.0, 6.9, widths=[6.8, 6.2], font_size=22, accent_col=1)
    add_box(s, 12.2, 10.35, 13.0, 2.0, fill=PALE_RED, line=RED, rounded=True)
    add_text(s, "点間移動・settling・光強度確率なし\n各時刻を独立な捕捉機会として評価", 12.7, 10.75, 12.0, 1.1, size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "平均捕捉時間は成功試行のみの条件付き平均。成功率を併記する。")

    s = new_content(prs, "非熱誤差は軌道・アライメント・姿勢・ドリフトを合成する", "PAT simulation conditions")
    add_box(s, 1.4, 3.1, 23.8, 2.1, fill=PALE, line=SKY, rounded=True)
    add_text(s, "enonthermal = eorbit + ealignment + eattitude + edrift", 2.0, 3.7, 22.6, 0.75, size=35, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_table(s, ["Component", "Model", "Timescale", "Size"], [("Orbit prediction", "Sentinel-1 TLE vs POEORB", "orbit-period family", "mean 280–700 µrad"), ("Alignment", "Gaussian constant / case", "DC", "50 µrad/axis, 1σ"), ("Attitude", "Gaussian / sample", "simplified broadband", "50 µrad/axis, 1σ"), ("Drift", "sinusoid", "900 s", "30 µrad amplitude")], 1.2, 6.0, 24.2, 5.0, widths=[4.8, 7.5, 5.4, 6.5], font_size=21, accent_col=0)
    add_text(s, "特定衛星の誤差バジェット再現ではなく、GNSS非搭載LEO小型衛星の共存シナリオ", 3.1, 11.75, 20.5, 0.65, size=27, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "軌道のみ実データ根拠が強く、他は代表値の簡易モデル。")

    s = new_content(prs, "利用可能な最新TLEをforward propagationし、PODと比較する", "Orbit prediction error")
    add_process(s, ["Sentinel-1\nNORAD 39634", "Latest TLE\nepoch≤t", "SGP4 forward\nrpred(t)", "AUX_POEORB\nrtruth(t)", "δr(t)"], 1.2, 3.1, 24.2)
    add_picture_fit(s, ROOT / "results/orbit/sentinel1_tle_vs_pod/orbit_prediction_error_3orbits.png", 1.2, 5.25, 15.4, 6.8)
    add_box(s, 17.2, 5.25, 8.0, 6.8, fill=PALE_GRAY, line=GRID, rounded=True)
    add_bullets(s, ["POEORBをtruth", "未来TLEのbackward伝搬なし", "60 s sampling", "相手機は片側truth", "TD時刻へcyclic resample"], 17.65, 5.8, 7.1, 4.8, size=23)
    add_text(s, "2026 S1データとTD軌道時刻は位相不一致", 17.7, 10.9, 7.0, 0.55, size=21, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "予測時刻より未来のTLEを使わない。現状は代表外乱としてcyclicに載せる。")

    s = new_content(prs, "位置誤差をリンクLOS横断面へ投影し、STT/body x–yへ変換する", "Orbit prediction error")
    add_box(s, 1.2, 3.1, 10.6, 8.3, fill=PALE_GRAY, line=GRID, rounded=True)
    add_text(s, "δr⊥ = δr − (δr·ℓ̂)ℓ̂", 1.8, 4.3, 9.4, 0.75, size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "θ ≈ δr⊥ / Rlink", 2.2, 5.8, 8.6, 0.75, size=32, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["orbit catalogから姿勢を構成", "STT/body x–yへ射影", "熱LOSと同じ2成分で加算"], 1.9, 7.4, 9.2, 2.7, size=24)
    add_table(s, ["Sun", "Partner / link", "Status"], [("MY", "along-track ISL 800 km", "ready"), ("PY", "anti-along ISL 800 km", "ready"), ("PX", "nadir ground ~695 km", "ready"), ("MX", "zenith proxy 800 km", "unrealistic")], 12.5, 3.1, 12.8, 6.9, widths=[2.8, 6.9, 3.1], font_size=21, accent_col=2)
    add_box(s, 12.5, 10.4, 12.8, 1.8, fill=PALE_RED, line=RED, rounded=True)
    add_text(s, "MXは角度計算可能だが現実的リンクではない", 13.0, 10.9, 11.8, 0.55, size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "以前の任意ECEF-Z基底からSTT/bodyへ接続。MX proxyに注意。")

    s = new_content(prs, "熱と軌道予測誤差はともに10⁻⁴ Hz台の軌道周期族となる", "Orbit prediction error")
    add_picture_fit(s, ROOT / "results/orbit/sentinel1_tle_vs_pod/orbit_error_stt_LTAN06_800km_1213COLD_MY_SUN_3orbits.png", 1.0, 3.0, 16.2, 8.8)
    add_table(s, ["Case", "mean norm"], [("MY ISL", "≈285 µrad"), ("PY ISL", "≈285 µrad"), ("PX ground", "≈694 µrad"), ("MX proxy", "≈612 µrad")], 17.9, 3.2, 7.2, 5.2, widths=[3.8, 3.4], font_size=23, accent_col=1)
    add_box(s, 17.9, 9.0, 7.2, 2.4, fill=PALE_ORANGE, line=ORANGE, rounded=True)
    add_text(s, "Thermal: ~101 min\nOrbit norm: ~101 / 50 min", 18.3, 9.45, 6.4, 1.2, size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "単純な『低周波＝熱』分離は難しい", 5.0, 12.1, 16.7, 0.6, size=29, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "熱も軌道も軌道位相依存。低周波残差を熱として取り込む説明は修正する。")

    s = new_content(prs, "熱のみでは階層モデルがthermal-truth上界にほぼ到達する", "PAT results")
    add_bars(s, ["No correction", "Static bias", "Hierarchical bcase", "Thermal truth"], [124.6, 4.78, 0.116, 0.100], 1.3, 3.5, 15.8, 6.5, colors=[RED, ORANGE, GREEN, BLUE], suffix=" s", max_value=125)
    add_table(s, ["Model", "Success", "Mean tacq"], [("No correction", "95.9%", "124.6 s"), ("Static bias", "99.3%", "4.78 s"), ("bcase", "100%", "0.116 s"), ("truth", "100%", "0.100 s")], 17.8, 3.3, 7.4, 5.4, widths=[3.4, 2.0, 2.0], font_size=20, accent_col=2)
    add_metric(s, "8.8 µrad", "mean thermal residual", 17.9, 9.2, 7.2, color=GREEN)
    add_text(s, "熱モデル自体の能力を見る理想条件。主結果は非熱込み。", 4.2, 11.9, 18.2, 0.6, size=25, color=GRAY, align=PP_ALIGN.CENTER)
    add_notes(s, "Staticが良いのはDC支配ケースが多いから。平均時間は成功試行条件付き。")

    s = new_content(prs, "非熱込みでも平均捕捉時間156.9→59.6 s、成功率94.0→97.1%となる", "PAT results")
    add_bars(s, ["No thermal correction", "Hierarchical bcase FF"], [156.9, 59.6], 1.5, 3.5, 15.2, 4.0, colors=[RED, GREEN], suffix=" s", max_value=160)
    add_metric(s, "−62%", "mean acquisition time", 17.6, 3.4, 7.3, color=GREEN)
    add_metric(s, "+3.1 pt", "success rate", 17.6, 5.75, 7.3, color=BLUE)
    add_table(s, ["Metric", "No correction", "bcase FF"], [("Success", "94.0%", "97.1%"), ("Mean tacq", "156.9 s", "59.6 s"), ("Median tacq", "150.3 s", "36.2 s"), ("Thermal residual", "667 µrad", "8.8 µrad")], 3.0, 8.35, 20.7, 3.55, widths=[8.1, 6.3, 6.3], font_size=21, accent_col=2)
    add_text(s, "補正後は非熱、特に軌道予測誤差が性能を支配する", 5.0, 12.25, 16.7, 0.6, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "21ケース単純平均。MY/PYで大幅改善、PX/MXは軌道誤差支配。")

    s = new_content(prs, "改善幅はリンク幾何と熱LOSの大きさに依存する", "PAT results")
    add_table(s, ["Case family", "Thermal LOS", "Orbit error", "PAT interpretation"], [("PY", "large (~1.2 mrad)", "~285 µrad", "large improvement"), ("MY", "150–250 µrad", "~285 µrad", "moderate improvement"), ("PX", "~0.9 mrad", "~694 µrad", "orbit remains large"), ("MX", "~0.9 mrad", "~612 µrad proxy", "interpret cautiously"), ("LTAN18 MY", "nearly DC ~150 µrad", "~290 µrad", "small thermal benefit")], 1.2, 3.1, 24.1, 7.0, widths=[4.2, 6.1, 5.5, 8.3], font_size=22, accent_col=3)
    add_box(s, 2.0, 10.75, 22.5, 1.7, fill=PALE_RED, line=RED, rounded=True)
    add_text(s, "全21ケース平均にはMXの非現実的zenith proxyが含まれる\n論文ではrealistic-link subsetも併記する候補", 2.6, 11.08, 21.3, 1.0, size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "全平均だけでなくケース依存性を見せる。MXをどう扱うか議論したい。")

    # Section 05
    new_section(prs, "05", "Discussion")

    s = new_content(prs, "Adaptiveは低周波熱抽出ではなく、主にモデルDC更新として考える", "Discussion")
    add_box(s, 1.2, 3.1, 11.8, 8.2, fill=PALE_GREEN, line=GREEN, rounded=True)
    add_text(s, "Proposed role", 1.7, 3.55, 10.8, 0.6, size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "b̂k+1 = b̂k + Kb rk", 2.1, 4.75, 10.0, 0.8, size=34, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["a：事前解析／熱真空試験で固定", "b_case：捕捉後残差で更新", "同一モード・複数パスで平均", "orbit/alignmentの吸収を制約"], 1.9, 6.2, 10.3, 4.0, size=25, accent=GREEN)
    add_box(s, 13.7, 3.1, 11.6, 8.2, fill=PALE_ORANGE, line=ORANGE, rounded=True)
    add_text(s, "Discussion points", 14.2, 3.55, 10.6, 0.6, size=28, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["a固定・b更新の分担", "同帯域の熱／軌道誤差の識別", "軌道誤差を別状態で同時推定", "現在のscan modelの妥当性", "ICSOまでのAdaptive実装範囲"], 14.4, 4.75, 10.2, 5.2, size=24, accent=ORANGE)
    add_notes(s, "ここから議論。周波数分離ではなく事前情報を使う状態推定が必要。")

    s = new_content(prs, "予測可能な熱成分は除去できたため、次はDC校正と非熱共存が課題となる", "Summary")
    add_table(s, ["What is established", "What remains"], [("TD/Femap → STT–LCT LOS → PAT", "real spacecraft validation"), ("shared |a|≈28–31 µrad/°C", "MZ / non-dominant axis"), ("LOO thermal residual ≈5.5 µrad", "power-continuous Level 2"), ("nonthermal tacq 156.9→59.6 s", "realistic-link subset / dynamics"), ("orbit and thermal share low-frequency band", "adaptive observability")], 1.3, 3.2, 24.0, 7.3, widths=[12.0, 12.0], font_size=23, accent_col=1)
    add_box(s, 2.0, 11.0, 22.5, 1.5, fill=PALE, line=SKY, rounded=True)
    add_text(s, "ICSOへ：モデル主張を固定し、PAT条件とAdaptiveの必要最小限を詰める", 2.6, 11.42, 21.3, 0.6, size=27, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "本編まとめ。ここで議論へ移る。")

    # Appendix section
    new_section(prs, "A", "Appendix")

    appendix = [
        ("関連研究の位置づけ", [("Badás et al. 2024", "FSOC opto-thermo-mechanical review"), ("Shi et al. 2023", "bus optimization → acquisition time"), ("Hu et al. 2022", "GEO Fourier compensation"), ("Li et al. 2025", "LEO geometry → NN LOS correction"), ("JANUS", "optical-head LOS≈KΔT"), ("Rüddenklau et al.", "body-pointing FF, not thermal")]),
        ("全解析ケースの設計", [("04,13–15", "MY", "ALL / +PROP / +PCDU / STTLCT"), ("08,16,18,19", "PY", "ALL / STTLCT / +PROP / +PCDU"), ("05,06,20,21", "PX", "ALL / STTLCT / +PROP / +PCDU"), ("09,17,23,24", "MX", "ALL / STTLCT / +PROP / +PCDU"), ("10–12", "MY", "HOT / Black / Alodine"), ("22", "MY", "PROP 12.5 W"), ("25", "MY", "LTAN18 / 693 km")]),
        ("熱光学特性", [("Alodine 1000", "α=0.150, ε=0.038, α/ε=3.947"), ("Black", "α=0.974, ε=0.920, α/ε=1.059"), ("0.5 synthetic", "α=0.500, ε=0.500")]),
        ("LOS定義の比較", []),
        ("Level-2係数と解釈", [("b0_MX", "+15.7 µrad"), ("b0_MY", "+2.8 µrad"), ("b0_PX", "−12.0 µrad"), ("b0_PY", "−24.0 µrad"), ("c_PROP", "−22.9 µrad"), ("c_PCDU", "−10.2 µrad")]),
        ("モデルの難しいケース", []),
        ("軌道予測誤差の仮定", [("Truth", "Sentinel-1 AUX_POEORB"), ("Prediction", "latest TLE + SGP4 forward"), ("Partner", "truth-known, one-sided error"), ("Range", "800 km ISL or nadir altitude"), ("Resampling", "cyclic to TD/Femap time")]),
        ("STT-frame射影結果", [("MY", "ISL along-track / 285 µrad"), ("PY", "ISL anti-along / 285 µrad"), ("PX", "ground nadir / 694 µrad"), ("MX", "zenith proxy / 612 µrad")]),
        ("PATスキャンの注意", [("Coverage", "40 µrad grid, 25 µrad radius"), ("Dynamics", "no slew / settling"), ("Target", "constant during one scan"), ("Statistics", "mean tacq conditional on success"), ("Opportunities", "each sample independent")]),
        ("Adaptive候補", [("Scalar EWMA", "b only; simplest"), ("Kalman filter", "b + orbit/alignment states"), ("Mode-conditioned table", "sun face × power mode"), ("Batch pass update", "avoid single-acquisition contamination")]),
        ("主要参考文献", []),
    ]
    for ai, (title, rows) in enumerate(appendix, 1):
        s = new_content(prs, title, "Appendix", page=f"A{ai}")
        if title == "LOS定義の比較":
            add_picture_fit(s, ROOT / "results/femap_deformation/15_LTAN06_800km_1213COLD_MY_STTLCT_HEAT_MY_0p5/los_definition_comparison.png", 1.2, 3.0, 24.0, 9.2)
        elif title == "モデルの難しいケース":
            for j, name in enumerate(["case11_bcase_true_vs_pred.png", "case22_bcase_true_vs_pred.png", "case25_bcase_true_vs_pred.png"]):
                add_picture_fit(s, ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/timeseries" / name, 0.8 + 8.45*j, 3.1, 7.9, 7.3)
            add_table(s, ["Case", "Issue"], [("11 Black", "oracle floor ≈13 µrad"), ("22 half power", "ON/OFF b mismatch"), ("25 all-sun", "thermal LOS nearly DC")], 3.0, 10.75, 20.7, 2.0, widths=[6.0, 14.7], font_size=21, accent_col=1)
        elif title == "主要参考文献":
            add_bullets(s, ["Riesing et al., TBIRD PAT on-orbit results, SPIE 2023", "Shi et al., Thermal Deformation Stability Optimization, 2023", "Badás et al., Opto-thermo-mechanical phenomena in satellite FSOC, 2024", "Hu et al., On-Board Thermal Motion Compensation, 2022", "Li et al., Thermal Deformation LOS Correction in LEO, 2025", "Rüddenklau et al., Feed-Forward Body-Pointing Compensation, 2024", "JANUS optical-head thermo-elastic LOS model"], 1.4, 3.1, 23.7, 8.7, size=25, accent=PURPLE)
        elif title == "全解析ケースの設計":
            add_table(s, ["Cases", "Sun", "Power / condition"], rows, 1.4, 3.1, 23.8, 8.0, widths=[4.5, 3.0, 16.3], font_size=23, accent_col=2, header=PURPLE)
        else:
            add_table(s, ["Item", "Detail"], rows, 2.0, 3.2, 22.6, min(8.5, 1.15*(len(rows)+1)), widths=[7.0, 15.6], font_size=23, accent_col=1, header=PURPLE)
        add_notes(s, f"Appendix: {title}")

    return prs


def main() -> None:
    prs = build()
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
