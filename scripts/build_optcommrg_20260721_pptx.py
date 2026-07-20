from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "papers/seminar/20260721_optcommrg_takamoto.pptx"
OUTPUT = ROOT / "papers/seminar/20260721_optcommrg_takamoto_updated.pptx"

NAVY = RGBColor(24, 49, 83)
BLUE = RGBColor(42, 104, 170)
TEAL = RGBColor(0, 142, 150)
CYAN = RGBColor(85, 190, 198)
ORANGE = RGBColor(231, 133, 48)
RED = RGBColor(190, 56, 64)
GREEN = RGBColor(49, 145, 97)
PURPLE = RGBColor(104, 79, 145)
INK = RGBColor(35, 43, 52)
MID = RGBColor(92, 103, 115)
LIGHT = RGBColor(238, 243, 247)
PALE_BLUE = RGBColor(230, 240, 250)
PALE_TEAL = RGBColor(225, 245, 244)
PALE_ORANGE = RGBColor(251, 239, 224)
PALE_RED = RGBColor(250, 232, 234)
WHITE = RGBColor(255, 255, 255)
GRID = RGBColor(207, 216, 224)

FONT_JP = "Yu Gothic"
FONT_LATIN = "Aptos"


def I(value: float):
    return Inches(value)


def clear_slides(prs: Presentation) -> None:
    for slide_id in list(prs.slides._sldIdLst):
        prs.slides._sldIdLst.remove(slide_id)


def set_fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor, width: float = 1.2) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 28,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.04,
    font: str = FONT_JP,
):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = I(margin)
    tf.margin_right = I(margin)
    tf.margin_top = I(margin)
    tf.margin_bottom = I(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_rich_text(
    slide,
    runs: Sequence[tuple[str, bool, RGBColor]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 28,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = I(0.04)
    tf.margin_top = tf.margin_bottom = I(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for text, bold, color in runs:
        r = p.add_run()
        r.text = text
        r.font.name = FONT_JP
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return shape


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = GRID,
    radius: bool = True,
    transparency: int = 0,
):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, I(x), I(y), I(w), I(h))
    set_fill(shape, fill, transparency)
    set_line(shape, line, 1.2)
    return shape


def add_label(slide, text: str, x: float, y: float, w: float, color: RGBColor = BLUE):
    shape = add_box(slide, x, y, w, 0.55, fill=color, line=color)
    add_text(slide, text, x, y + 0.02, w, 0.46, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_bullets(
    slide,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 28,
    color: RGBColor = INK,
    bullet_color: RGBColor = TEAL,
    gap: float = 8,
):
    shape = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = I(0.05)
    tf.margin_right = I(0.04)
    tf.margin_top = I(0.03)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        r0 = p.add_run()
        r0.text = "●  "
        r0.font.name = FONT_JP
        r0.font.size = Pt(size - 4)
        r0.font.color.rgb = bullet_color
        r1 = p.add_run()
        r1.text = item
        r1.font.name = FONT_JP
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
    return shape


def add_title(slide, title: str, section: str, number: int, total: int, *, appendix: bool = False):
    add_text(slide, section, 0.6, 0.22, 15.5, 0.55, size=20, color=TEAL, bold=True)
    add_text(slide, title, 1.15, 1.15, 24.3, 1.1, size=38, color=NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, I(1.15), I(2.23), I(24.3), I(0.055))
    set_fill(line, TEAL if not appendix else PURPLE)
    line.line.fill.background()
    add_text(slide, f"{'A' if appendix else ''}{number}", 24.2, 14.25, 1.65, 0.42, size=18, color=MID, align=PP_ALIGN.RIGHT)


def new_slide(prs: Presentation, title: str, section: str, *, appendix: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    number = len(prs.slides) - 35 if appendix else len(prs.slides)
    add_title(slide, title, section, number, 0, appendix=appendix)
    return slide


def add_notes(slide, notes: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = notes.strip()


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = TEAL, width: float = 3):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def add_picture_fit(slide, path: Path, x: float, y: float, w: float, h: float, *, border: bool = True):
    if not path.exists():
        box = add_box(slide, x, y, w, h, fill=PALE_RED, line=RED)
        add_text(slide, f"Missing image:\n{path.name}", x + 0.2, y + 0.2, w - 0.4, h - 0.4, size=22, color=RED)
        return box
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    if border:
        frame = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, I(x), I(y), I(w), I(h)
        )
        frame.fill.background()
        set_line(frame, GRID, 1.0)
    pic = slide.shapes.add_picture(str(path), I(px), I(py), I(pw), I(ph))
    return pic


def add_table(
    slide,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    widths: Sequence[float] | None = None,
    font_size: float = 22,
    header_color: RGBColor = NAVY,
    accent_col: int | None = None,
):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), I(x), I(y), I(w), I(h))
    table = shape.table
    if widths:
        for col, cw in zip(table.columns, widths):
            col.width = I(cw)
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
    for r, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            if accent_col == c:
                cell.fill.fore_color.rgb = PALE_TEAL
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
    for r in range(len(rows) + 1):
        for c in range(len(headers)):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = I(0.08)
            cell.margin_top = cell.margin_bottom = I(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_JP
                    run.font.size = Pt(font_size)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else INK
    return shape


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, color: RGBColor = TEAL):
    add_box(slide, x, y, w, 1.65, fill=WHITE, line=color)
    add_text(slide, value, x + 0.1, y + 0.12, w - 0.2, 0.75, size=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, y + 0.91, w - 0.2, 0.5, size=19, color=MID, align=PP_ALIGN.CENTER)


def add_process(slide, labels: Sequence[str], x: float, y: float, total_w: float, *, colors: Sequence[RGBColor] | None = None):
    n = len(labels)
    gap = 0.55
    bw = (total_w - gap * (n - 1)) / n
    colors = colors or [NAVY, BLUE, TEAL, GREEN, ORANGE, PURPLE]
    for idx, label in enumerate(labels):
        bx = x + idx * (bw + gap)
        add_box(slide, bx, y, bw, 1.25, fill=colors[idx % len(colors)], line=colors[idx % len(colors)])
        add_text(slide, label, bx + 0.08, y + 0.12, bw - 0.16, 0.95, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if idx < n - 1:
            add_arrow(slide, bx + bw + 0.08, y + 0.63, bx + bw + gap - 0.08, y + 0.63, color=MID, width=2)


def add_horizontal_bars(slide, labels, values, x, y, w, h, *, colors=None, suffix="", max_value=None):
    colors = colors or [BLUE] * len(values)
    max_value = max_value or max(values)
    row_h = h / len(values)
    label_w = w * 0.30
    value_w = w - label_w - 0.9
    for i, (label, value, color) in enumerate(zip(labels, values, colors)):
        yy = y + i * row_h
        add_text(slide, label, x, yy + 0.06, label_w - 0.1, row_h - 0.12, size=20, color=INK, valign=MSO_ANCHOR.MIDDLE)
        add_box(slide, x + label_w, yy + 0.15, value_w, row_h - 0.3, fill=LIGHT, line=LIGHT, radius=False)
        bw = max(0.08, value_w * value / max_value)
        add_box(slide, x + label_w, yy + 0.15, bw, row_h - 0.3, fill=color, line=color, radius=False)
        add_text(slide, f"{value:g}{suffix}", x + label_w + value_w + 0.12, yy + 0.05, 0.75, row_h - 0.1, size=20, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)


def title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, I(0), I(0), I(8.6), I(15))
    set_fill(band, NAVY); band.line.fill.background()
    add_text(slide, "OPTICAL COMMUNICATION\nRESEARCH GROUP", 1.05, 1.15, 6.5, 1.7, size=26, color=CYAN, bold=True)
    add_text(slide, "2026 / 07 / 21", 1.05, 12.4, 5.0, 0.6, size=24, color=WHITE, bold=True)
    add_text(slide, "Hideki Takamoto\nThe University of Tokyo", 1.05, 13.15, 6.3, 1.0, size=22, color=WHITE)
    add_text(slide, "Feedforward and Adaptive Correction of\nTime-Varying Thermal Bias for Coarse Acquisition\nin Optical Communication Systems", 9.6, 3.0, 15.5, 3.4, size=38, color=NAVY, bold=True)
    add_text(slide, "熱ひずみ予測と適応補正による\n光通信粗捕捉性能向上の検討", 9.65, 7.2, 14.8, 2.0, size=34, color=TEAL, bold=True)
    add_text(slide, "Progress update: hierarchical ΔT model, PAT evaluation,\nand orbit-prediction-error integration", 9.65, 10.2, 14.8, 1.2, size=25, color=MID)
    add_notes(slide, "本日は前回7月6日の発表以降の進捗を中心に話す。特に、21ケースを共通係数で説明する階層ΔTモデル、TLE–POD軌道予測誤差を含めたPAT評価、Adaptiveを何に使うべきかを議論したい。")


def build_deck() -> Presentation:
    source = Presentation(str(SOURCE))
    prs = Presentation()
    prs.slide_width = source.slide_width
    prs.slide_height = source.slide_height
    title_slide(prs)

    # 2
    s = new_slide(prs, "今日のゴール：モデルを見せ、仮定をレビューしてもらう", "Overview")
    add_process(s, ["Thermo-structural\nanalysis", "Hierarchical\nΔT model", "PAT + orbit\nerror", "Adaptive\ndiscussion"], 1.5, 3.0, 23.2)
    add_box(s, 1.5, 5.2, 23.2, 5.9, fill=LIGHT, line=GRID)
    add_text(s, "今日特に議論したいこと", 2.0, 5.65, 8.0, 0.7, size=30, color=NAVY, bold=True)
    add_bullets(s, [
        "温度差感度 a を事前同定し、ケースDC b_case を軌道上更新する分担は妥当か",
        "熱と軌道予測誤差が同じ軌道周期帯にあるとき、何を観測可能と考えるべきか",
        "現在の粗捕捉スキャン・誤差バジェットは、通信性能評価として十分か",
    ], 2.0, 6.5, 21.7, 3.8, size=29)
    add_notes(s, "今日は完成報告ではなく設計レビューに近い。前半で解析条件を再確認し、後半で軌道誤差とAdaptiveの論点に時間を残す。")

    # 3
    s = new_slide(prs, "熱ひずみは粗捕捉開始時の scan-center error として残る", "Background")
    add_box(s, 1.4, 3.0, 7.0, 7.6, fill=PALE_BLUE, line=BLUE)
    add_text(s, "光フィードバック前", 1.8, 3.45, 6.2, 0.8, size=31, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "姿勢・軌道・アライメント・\n熱LOSを含む不確定領域を探索", 1.9, 5.0, 6.0, 1.8, size=28, align=PP_ALIGN.CENTER)
    add_text(s, "SCAN", 3.0, 8.0, 3.8, 1.0, size=42, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(s, 8.8, 6.8, 11.0, 6.8, color=TEAL, width=4)
    add_box(s, 11.3, 3.0, 6.0, 7.6, fill=PALE_ORANGE, line=ORANGE)
    add_text(s, "日照・蝕／機器発熱", 11.7, 3.45, 5.2, 0.8, size=29, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "温度場 → 差動膨張 →\nSTT–LCT相対回転", 11.8, 5.0, 5.0, 1.8, size=28, align=PP_ALIGN.CENTER)
    add_text(s, "10²–10³ µrad", 12.1, 8.0, 4.4, 0.9, size=34, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(s, 17.7, 6.8, 19.9, 6.8, color=TEAL, width=4)
    add_box(s, 20.1, 3.0, 5.1, 7.6, fill=PALE_TEAL, line=TEAL)
    add_text(s, "本研究", 20.6, 3.45, 4.1, 0.8, size=31, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "予測可能な熱成分を\nscan centerから事前に引く", 20.55, 5.0, 4.2, 2.2, size=28, align=PP_ALIGN.CENTER)
    add_text(s, "探索域・捕捉時間↓", 20.5, 8.1, 4.3, 0.9, size=29, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "構造を変えず、運用側から粗捕捉性能を改善する", 5.0, 12.0, 17.0, 0.8, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "光RGには基本説明なので短く。熱変形そのものではなく、STT基準とLCT光軸の差がscan center errorになる点だけ再確認する。")

    # 4
    s = new_slide(prs, "前回7/6：熱補正が効くことは確認、軽量モデルは未確立", "Progress since 7/6")
    add_label(s, "7/6 preliminary", 1.4, 3.0, 5.2, ORANGE)
    add_bullets(s, ["TD/Femap LOSをPATへ接続", "熱真値に近いFF補正で捕捉改善", "非熱誤差は簡易仮定", "Adaptiveは概念ブロックのみ"], 1.5, 3.9, 10.1, 5.8, size=28, bullet_color=ORANGE)
    add_box(s, 12.4, 3.0, 12.8, 7.7, fill=LIGHT, line=GRID)
    add_text(s, "当時の予備結果", 12.9, 3.55, 5.0, 0.7, size=28, color=NAVY, bold=True)
    add_metric(s, "25.8 s", "mean acquisition time", 13.0, 4.8, 5.0, ORANGE)
    add_metric(s, "−70%", "vs. no correction", 19.1, 4.8, 5.0, TEAL)
    add_text(s, "注意：ケース・非熱条件・集計方法が現在と異なるため、今回結果とは直接比較しない", 13.0, 7.25, 11.2, 1.6, size=25, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "今回の問い：truth correctionではなく、運用可能な軽量モデルでどこまで近づけるか", 3.0, 11.7, 20.7, 1.0, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "25.8秒という数字を更新前の予備結果として位置づける。今回の評価系はケース集合も非熱誤差も変わっているので、改善率の直接比較をしない。")

    # 5
    s = new_slide(prs, "2週間で『理想補正』から『ケース横断モデル＋現実的外乱』へ", "Progress since 7/6")
    rows = [
        ("熱LOSモデル", "thermal truth / Fourier候補", "階層 ΔT：共有 a + b_case"),
        ("解析ケース", "数ケース", "21ケース（太陽面・発熱・被覆・軌道）"),
        ("検証", "within-case中心", "LOO b + 後続2軌道test"),
        ("軌道誤差", "Gaussian placeholder", "Sentinel-1 TLE vs POEORB"),
        ("座標系", "任意LOS横断2軸", "orbit catalogからSTT/bodyへ射影"),
        ("Adaptive", "低周波残差を学習", "主にb_case/DC更新へ再定義"),
    ]
    add_table(s, ["項目", "7/6", "7/21"], rows, 1.35, 3.05, 23.9, 7.9, widths=[4.2, 8.4, 11.3], font_size=24, accent_col=2)
    add_text(s, "現在の主張：固定少数係数＋温度差で熱主成分を落とし、同帯域の非熱誤差と共存しても捕捉を改善できる", 2.0, 11.75, 22.6, 1.4, size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "ここが進捗の要約。以降はこの表の各項目を順に説明する。")

    # 6
    s = new_slide(prs, "解析からPATまでの一気通貫パイプライン", "Thermo-structural analysis")
    add_process(s, ["Orbit / attitude\nheat & coating", "Thermal Desktop\ntemperature", "Femap / Nastran\nrotation", "STT–LCT\nthermal LOS", "Lightweight\nmodel", "PAT\nscan center"], 1.1, 3.2, 24.5)
    add_box(s, 1.6, 5.5, 10.8, 5.4, fill=PALE_BLUE, line=BLUE)
    add_text(s, "High-fidelity reference", 2.0, 5.9, 10.0, 0.7, size=29, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["軌道熱環境・日照/蝕", "温度場・熱変形", "STTとLCTの相対回転"], 2.2, 6.9, 9.6, 3.2, size=27, bullet_color=BLUE)
    add_box(s, 14.0, 5.5, 10.8, 5.4, fill=PALE_TEAL, line=TEAL)
    add_text(s, "Operational layer", 14.4, 5.9, 10.0, 0.7, size=29, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["温度センサ＋運用フラグ", "予測LOSをscan centerへ", "捕捉後残差でモデル更新"], 14.6, 6.9, 9.6, 3.2, size=27)
    add_text(s, "TD/Femapはオンボード実装ではなく、軽量モデルの教師・評価用真値", 4.5, 12.0, 17.8, 0.8, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "解析パイプライン全体。今回のhigh-fidelityは設計用であり、オンボードではΔTモデルだけを使う。")

    # 7
    s = new_slide(prs, "衛星モデル・コンポ配置・座標系", "Thermo-structural analysis")
    add_box(s, 1.4, 3.0, 12.0, 8.6, fill=LIGHT, line=GRID)
    add_text(s, "Baseline bus", 1.9, 3.35, 4.0, 0.7, size=29, color=NAVY, bold=True)
    # satellite sketch
    add_box(s, 4.0, 5.0, 6.3, 4.3, fill=PALE_BLUE, line=BLUE, radius=False)
    add_text(s, "PZ: STT", 5.6, 4.15, 3.1, 0.6, size=25, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "MZ: LCT\nboresight ≈ −Z", 5.2, 6.1, 4.0, 1.2, size=26, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "MY: PCDU", 1.9, 6.0, 2.1, 0.6, size=23, color=ORANGE, bold=True)
    add_text(s, "PY: PROP", 10.4, 6.0, 2.2, 0.6, size=23, color=ORANGE, bold=True)
    add_arrow(s, 7.1, 8.0, 7.1, 10.2, color=TEAL, width=4)
    add_text(s, "LCT LOS", 5.7, 10.15, 2.8, 0.6, size=23, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_table(s, ["項目", "値"], [("外形", "590 × 600 × 990 mm"), ("パネル", "10 mm shell"), ("材料", "A5052"), ("LCT", "MZ center"), ("STT", "PZ center")], 14.2, 3.2, 10.9, 7.2, widths=[3.8, 7.1], font_size=24)
    add_text(s, "太陽面 MX / MY / PX / PY によって、熱変形の支配軸と符号が変わる", 3.2, 12.15, 20.4, 0.8, size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "箱型の簡略衛星モデル。LCTはMZ、STTはPZ、発熱コンポは±Y側に配置。支配軸と太陽面の対応を見るため座標を確認する。")

    # 8
    s = new_slide(prs, "PAT入力には STT基準で見たLCT光軸の相対回転を用いる", "Thermo-structural analysis")
    img = ROOT / "results/femap_deformation/15_LTAN06_800km_1213COLD_MY_STTLCT_HEAT_MY_0p5/los_definition_comparison.png"
    add_picture_fit(s, img, 1.3, 3.0, 14.0, 8.8)
    add_box(s, 16.0, 3.0, 9.1, 8.8, fill=LIGHT, line=GRID)
    add_text(s, "採用：far_field_los", 16.5, 3.5, 8.1, 0.8, size=30, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, [
        "LCT光軸回転 − STT姿勢基準回転",
        "遠方通信でscan centerへ効く角度",
        "代表点間並進によるcenterline tiltは加えない",
        "stt_relative_losは角度バジェット確認用",
    ], 16.55, 4.7, 8.0, 4.5, size=25)
    add_text(s, "θthermal,true(t) = [θx(t), θy(t)]", 16.8, 10.1, 7.5, 0.7, size=27, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "LOS定義の混乱を避ける。PATの主入力はfar-field relative rotation。並進由来のcenterline tiltを重ねた量は診断用に残している。")

    # 9
    s = new_slide(prs, "軌道熱解析は基本3軌道・約60秒刻み", "Thermo-structural analysis")
    rows = [("Baseline orbit", "LTAN06 / 800 km / COLD"), ("Orbit period", "≈ 6050 s (101 min)"), ("Duration", "≈ 18157 s (3 orbits)"), ("Sampling", "≈ 60.5 s / 301 samples"), ("Eclipse history", "TD LOGIC_SUN"), ("Sun-facing panels", "MX / MY / PX / PY")]
    add_table(s, ["Parameter", "Setting"], rows, 1.4, 3.0, 11.4, 7.9, widths=[4.2, 7.2], font_size=24)
    add_box(s, 13.7, 3.0, 11.4, 7.9, fill=PALE_BLUE, line=BLUE)
    add_text(s, "Additional environments", 14.2, 3.45, 10.4, 0.7, size=29, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["HOT thermal environment", "Black / Alodine / synthetic α=ε=0.5", "LTAN18 / 693 km Sentinel-1 proxy", "case25: nearly all-sun, thermal LOS nearly DC"], 14.4, 4.6, 10.0, 4.5, size=26, bullet_color=BLUE)
    add_text(s, "蝕入り・蝕明けは時刻表ではなく、TDが出力したLOGIC_SUN時系列を正本とする", 3.2, 12.0, 20.4, 0.9, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "基本は3軌道。モデル同定に最初の1軌道、残り2軌道を時系列testに使う。case25は異なる軌道条件の確認。")

    # 10
    s = new_slide(prs, "Femap構造解析条件：簡略shellモデルを高忠実度参照として用いる", "Thermo-structural analysis")
    rows = [("Material", "Aluminum 5052 Annealed Wrought"), ("Young's modulus", "70.327 GPa"), ("Poisson ratio", "0.33"), ("CTE", "2.376×10⁻⁵ /°C"), ("Density", "2685 kg/m³"), ("Reference temperature", "23.9 °C"), ("Constraint", "small region near STT reference")]
    add_table(s, ["Structural parameter", "Value"], rows, 1.4, 3.0, 11.5, 8.7, widths=[5.2, 6.3], font_size=23)
    img = ROOT / "results/femap_deformation/15_LTAN06_800km_1213COLD_MY_STTLCT_HEAT_MY_0p5/stt_lct_motion_overview.png"
    add_picture_fit(s, img, 13.7, 3.0, 11.5, 8.7)
    add_text(s, "TD mapper温度 → 各時刻の熱荷重 → STT/LCT基準面の変位・回転", 3.5, 12.3, 19.7, 0.8, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "高忠実度という言葉は相対的。実機詳細構造ではなく、熱変形メカニズムと軽量モデル評価の参照となる簡略shellモデル。拘束条件依存性も残る。")

    # 11
    s = new_slide(prs, "内部発熱モードで ΔT とケースDCの両方を変化させる", "Thermo-structural analysis")
    add_table(s, ["Component", "Nominal heat", "Location"], [("LCT", "10 W", "MZ center"), ("STT", "1.5 W", "PZ center"), ("PROP", "25 W", "PY side"), ("PCDU", "10 W", "MY side")], 1.4, 3.0, 10.2, 5.8, widths=[3.3, 3.0, 3.9], font_size=24)
    add_table(s, ["Power mode", "LCT/STT", "PROP", "PCDU"], [("STTLCT", "ON", "OFF", "OFF"), ("+PROP", "ON", "25 W", "OFF"), ("+PCDU", "ON", "OFF", "10 W"), ("ALL", "ON", "25 W", "10 W"), ("Half-power test", "ON", "12.5 W", "10 W")], 12.2, 3.0, 13.0, 6.7, widths=[4.6, 2.8, 2.8, 2.8], font_size=22, accent_col=0)
    add_box(s, 2.3, 10.4, 21.8, 2.2, fill=PALE_ORANGE, line=ORANGE)
    add_rich_text(s, [("発熱の主効果", True, ORANGE), ("は面間温度差 ΔT に入る一方、", False, INK), ("残りのDC", True, RED), ("が b_case に現れる", False, INK)], 3.0, 10.95, 20.4, 0.9, size=31, align=PP_ALIGN.CENTER)
    add_notes(s, "発熱モードはLevel 2の説明変数になる。PROP/PCDUを足した効果の多くはΔTへ吸収され、残差DCだけがbに残る。")

    # 12
    s = new_slide(prs, "21ケースで太陽面・発熱・被覆・軌道の感度を横断評価", "Thermo-structural analysis")
    rows = [("Sun face", "MX / MY / PX / PY", "支配軸・符号"), ("Power mode", "STTLCT / +PROP / +PCDU / ALL", "ΔT・DC"), ("Coating", "0.5 / Black / Alodine", "時変振幅・残差床"), ("Environment", "COLD / HOT", "ケースDC"), ("Power level", "PROP 25 / 12.5 W", "ON/OFFモデル限界"), ("Orbit", "LTAN06 / LTAN18", "別熱履歴")]
    add_table(s, ["Varied factor", "Cases", "What it tests"], rows, 1.35, 3.0, 23.9, 7.7, widths=[4.4, 9.0, 10.5], font_size=24, accent_col=2)
    add_box(s, 2.0, 11.3, 22.6, 1.5, fill=PALE_RED, line=RED)
    add_text(s, "除外：01–03, 07（MZ太陽指向）— 現在の1軸sunfaceモデルの対象外", 2.4, 11.7, 21.8, 0.6, size=27, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "ケース設計の全体像。21ケースはいずれも基本構造が同じ感度解析で、独立機体への汎化試験ではない。")

    # 13
    s = new_slide(prs, "ケースIDを軸に解析条件と成果物を一貫管理", "Thermo-structural analysis")
    add_process(s, ["case_matrix.xlsx\norbit_catalog.xlsx", "TD temperature\nLOGIC_SUN", "mapper\noutput.dat", "Femap\nrotation", "los_angles.csv", "lightweight\ndataset"], 1.0, 3.2, 24.7)
    add_box(s, 1.5, 5.5, 23.7, 5.4, fill=LIGHT, line=GRID)
    add_text(s, "自動化・再現性のポイント", 2.1, 5.95, 7.0, 0.7, size=29, color=NAVY, bold=True)
    add_bullets(s, [
        "case_idから下流パスを規則解決し、Excelの手書きパス列を廃止",
        "TD温度場の妥当性チェック後にFemap・LOSへ流す",
        "同じcase_idで温度・LOS・モデル・PAT結果を追跡",
        "cases 4–25を同一コマンドで再評価可能",
    ], 2.2, 6.9, 22.0, 3.2, size=26)
    add_text(s, "解析数を増やすことより、同じ条件を再現し比較できることを優先", 4.0, 12.0, 18.7, 0.8, size=29, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "ケース増加で最も大きかった作業は自動化と妥当性チェック。過去に光学特性未反映やmapper混線を経験したため、ケースID管理を強化した。")

    # 14
    s = new_slide(prs, "代表ケース：面間温度差と熱LOSが同じ軌道周期で変動", "Thermo-structural results")
    img1 = ROOT / "results/femap_deformation/04_LTAN06_800km_1213COLD_MY_ALL_HEAT_MY_0p5/default_surface_9points_temperature_overview.png"
    img2 = ROOT / "results/femap_deformation/04_LTAN06_800km_1213COLD_MY_ALL_HEAT_MY_0p5_far_field_los_angle_budget.png"
    add_picture_fit(s, img1, 1.25, 3.0, 11.9, 8.3)
    add_picture_fit(s, img2, 13.5, 3.0, 11.9, 8.3)
    add_label(s, "Temperature", 4.6, 11.55, 4.8, BLUE)
    add_label(s, "Far-field LOS", 17.0, 11.55, 4.8, TEAL)
    add_text(s, "Case 04: LTAN06 / MY sun / ALL heat / MY α=ε=0.5", 4.0, 12.55, 18.7, 0.6, size=24, color=MID, align=PP_ALIGN.CENTER)
    add_notes(s, "代表ケース04。温度時系列とLOS時系列を対応づけ、太陽面と反対面の温度差が支配軸LOSを説明できそうだと分かった。")

    # 15
    s = new_slide(prs, "太陽面で支配軸と符号が決まり、感度の絶対値はほぼ共通", "Thermo-structural results")
    axes = [("MX", "X", "+", 30.6, PALE_TEAL), ("PX", "X", "−", 28.1, PALE_BLUE), ("MY", "Y", "+", 28.6, PALE_ORANGE), ("PY", "Y", "−", 28.7, PALE_RED)]
    for i, (face, axis, sign, val, fill) in enumerate(axes):
        x = 1.4 + i * 6.0
        add_box(s, x, 3.2, 5.3, 6.8, fill=fill, line=TEAL if sign == "+" else BLUE)
        add_text(s, face, x + 0.2, 3.65, 4.9, 0.8, size=37, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, f"dominant {axis}-axis", x + 0.25, 5.0, 4.8, 0.7, size=25, color=MID, align=PP_ALIGN.CENTER)
        add_text(s, sign, x + 1.3, 6.0, 2.7, 1.2, size=54, color=TEAL if sign == "+" else BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, f"|a| ≈ {val:.1f}\nµrad/°C", x + 0.4, 7.55, 4.5, 1.5, size=28, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "対向面で符号反転、MX/PXはX軸、MY/PYはY軸支配 — 構造対称性と整合", 2.2, 11.25, 22.2, 1.0, size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "ここでは後で示す共有aの結果を物理発見として先に見せる。太陽面の対称性と符号反転がモデルの根拠。")

    # 16
    s = new_slide(prs, "被覆・発熱・軌道条件は振幅とケースDCを異なる形で変える", "Thermo-structural results")
    add_table(s, ["Factor", "Observed effect", "Model implication"], [
        ("Coating", "Blackで時変振幅・残差床が増加", "Level 1の限界"),
        ("PROP/PCDU heat", "平均値とΔTが変化", "ΔT + b_case"),
        ("HOT", "時変が小さくほぼDC", "Level 2の環境拡張"),
        ("LTAN18 all-sun", "熱LOSがほぼ一定", "Staticでも十分な場合"),
        ("Half PROP power", "aは不変、bが中間値", "0/1フラグの限界"),
    ], 1.4, 3.0, 23.8, 7.5, widths=[5.0, 9.1, 9.7], font_size=24, accent_col=2)
    add_box(s, 2.2, 11.15, 22.2, 1.7, fill=PALE_TEAL, line=TEAL)
    add_text(s, "時変とケース間DCを同じ回帰に押し込まず、階層分離する必要がある", 2.7, 11.65, 21.2, 0.7, size=30, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "感度解析の結論。時系列の時変と運用モード間のDCを分けることが階層モデルの動機。")

    # 17
    s = new_slide(prs, "軌道上モデルは少数入力・解釈可能・PATへ直結できることが必要", "Lightweight LOS model")
    add_box(s, 1.4, 3.0, 8.0, 7.8, fill=PALE_BLUE, line=BLUE)
    add_text(s, "Operational requirements", 1.9, 3.45, 7.0, 0.7, size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["TD/Femapをオンボードで回さない", "少数温度センサ", "太陽面・運用モード変更", "固定少数係数", "scan-center correctionを直接出力"], 2.0, 4.5, 6.8, 5.2, size=25, bullet_color=BLUE)
    add_table(s, ["Candidate", "Pros / cons"], [("Static bias", "軽量 / 時変を追えない"), ("Fourier", "周期を表現 / 軌道条件依存"), ("Generic temperature regression", "柔軟 / 共線・解釈性"), ("Sunface ΔT", "物理的・少数入力"), ("Hierarchical ΔT", "時変＋運用DC")], 10.2, 3.0, 15.0, 7.8, widths=[5.5, 9.5], font_size=24, accent_col=0)
    add_text(s, "本命：太陽面−反対面温度差 ΔT と運用フラグ", 4.2, 11.75, 18.2, 0.8, size=31, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Fourierも候補だったが、物理入力で軌道条件の変化を扱いたい。ΔTモデルは温度センサと運用情報で実装可能。")

    # 18
    s = new_slide(prs, "失敗：コンポ温度を時系列特徴に足すと、共線と低SNRで係数が壊れる", "Lightweight LOS model")
    add_box(s, 1.4, 3.0, 11.3, 7.7, fill=PALE_RED, line=RED)
    add_text(s, "Attempt", 1.9, 3.45, 3.0, 0.7, size=29, color=RED, bold=True)
    add_text(s, "LOS = b + a·ΔT\n      + c·(Tattach−Tref)", 2.3, 4.7, 9.4, 1.8, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["corr(ΔT, TPCDU) ≈ 0.995", "局所差の軌道内std ≈ 0.1°C", "RMSEは下がっても係数が巨大化", "ケース間で係数が不安定"], 2.0, 7.1, 9.8, 3.0, size=24, bullet_color=RED)
    add_arrow(s, 13.2, 6.8, 15.1, 6.8, color=TEAL, width=4)
    add_box(s, 15.4, 3.0, 9.8, 7.7, fill=PALE_TEAL, line=TEAL)
    add_text(s, "Interpretation", 15.9, 3.45, 8.8, 0.7, size=29, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "効いていたのは\n軌道内の微小変動ではなく\nケース平均DC", 16.3, 5.0, 8.0, 2.4, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "→ コンポ効果を b_case 側へ", 16.3, 8.25, 8.0, 0.9, size=29, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "モデル選定はRMSEだけでなく、係数の安定性と運用解釈で判断", 3.6, 11.75, 19.5, 0.8, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "失敗結果を明示する。コンポ温度時系列を入れると見かけのRMSEは改善するが、ΔTとの共線で係数が解釈不能になった。これが階層化の直接の動機。")

    # 19
    s = new_slide(prs, "Level 1：軌道内の時変は面間温度差一本で表す", "Hierarchical ΔT model")
    add_box(s, 1.6, 3.2, 23.4, 3.1, fill=PALE_TEAL, line=TEAL)
    add_text(s, "θdom(t)  ≈  bcase  +  asunface · ΔT(t)", 2.1, 4.05, 22.4, 1.2, size=43, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "ΔT(t) = Tsunface(t) − Topposite(t)", 5.0, 6.85, 16.7, 0.9, size=32, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_box(s, 1.8, 8.4, 7.0, 3.2, fill=PALE_BLUE, line=BLUE)
    add_text(s, "a(sun)", 2.3, 8.8, 6.0, 0.7, size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "差動膨張・曲げ感度\n面ごとに共有", 2.4, 9.7, 5.8, 1.2, size=26, align=PP_ALIGN.CENTER)
    add_box(s, 9.8, 8.4, 7.0, 3.2, fill=PALE_ORANGE, line=ORANGE)
    add_text(s, "bcase", 10.3, 8.8, 6.0, 0.7, size=30, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "ΔTに入りきらない\nケース内定数", 10.4, 9.7, 5.8, 1.2, size=26, align=PP_ALIGN.CENTER)
    add_box(s, 17.8, 8.4, 7.0, 3.2, fill=LIGHT, line=GRID)
    add_text(s, "dominant axis", 18.3, 8.8, 6.0, 0.7, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "MX/PX→X\nMY/PY→Y", 18.4, 9.7, 5.8, 1.2, size=26, align=PP_ALIGN.CENTER)
    add_notes(s, "Level 1は支配軸だけを対象とする。非支配軸は現状モデル外。aは構造感度、bはケースDC。")

    # 20
    s = new_slide(prs, "Level 2：ケースDCを太陽面と発熱モードで説明する", "Hierarchical ΔT model")
    add_box(s, 1.6, 3.2, 23.4, 3.1, fill=PALE_ORANGE, line=ORANGE)
    add_text(s, "bcase  ≈  b0(sunface) + cprop Iprop + cpcdu Ipcdu", 2.0, 4.05, 22.6, 1.2, size=40, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_table(s, ["Case-dependent input", "Cross-case fixed parameter"], [("ΔT(t) timeseries", "a(MX/MY/PX/PY) × 4"), ("sun_face", "b0(MX/MY/PX/PY) × 4"), ("PROP/PCDU ON/OFF", "cprop, cpcdu × 2")], 2.0, 7.0, 22.6, 4.2, widths=[11.3, 11.3], font_size=25, accent_col=1)
    add_text(s, "合計10個のスカラー係数で21ケースを表現", 5.2, 12.05, 16.3, 0.8, size=32, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Level 2は運用モード間のDCを予測する。現在の実装では発熱フラグをMY/PY太陽面で有効とし、MX/PXでは差が小さいため無効化。")

    # 21
    s = new_slide(prs, "係数同定：先頭1軌道で時変、ケース横断LOOでDCを評価", "Hierarchical ΔT model")
    add_process(s, ["Per case\nfirst orbit", "Fit aemp, bemp\nRidge ≈ OLS", "Median aemp\nper sun face", "Level-2 fit\nacross cases", "LOO bpred\n+ next 2 orbits"], 1.2, 3.15, 24.2)
    add_box(s, 1.5, 5.55, 23.7, 5.7, fill=LIGHT, line=GRID)
    add_text(s, "評価の切り分け", 2.0, 5.95, 5.0, 0.7, size=29, color=NAVY, bold=True)
    add_table(s, ["Layer", "Train", "Test / validation"], [("Level 1", "各ケース先頭1軌道", "残り2軌道"), ("Level 2 b", "他20ケース", "Leave-one-case-out"), ("Shared a", "面ごとの全aemp中央値", "完全LOOではない")], 2.0, 6.9, 22.6, 3.5, widths=[4.2, 8.5, 9.9], font_size=23, accent_col=2)
    add_text(s, "21ケースは同一基本構造の感度ケース：別機体への汎化を示すものではない", 3.2, 11.85, 20.4, 0.8, size=27, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "過大評価を避けるため評価分割を明示。bはLOOだがa_sharedは対象ケースを含む中央値で、完全LOOではない。ただしケース追加後もほぼ不変。")

    # 22
    s = new_slide(prs, "共有感度 a は 21ケース追加後も ±28–31 µrad/°Cで安定", "Hierarchical ΔT model")
    img = ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/bcase_a_emp_by_sunface.png"
    add_picture_fit(s, img, 1.25, 3.0, 16.0, 8.6)
    add_table(s, ["Sun", "a_shared"], [("MX", "+30.6"), ("MY", "+28.6"), ("PX", "−28.1"), ("PY", "−28.7")], 18.0, 3.2, 7.0, 5.5, widths=[3.0, 4.0], font_size=25, accent_col=1)
    add_box(s, 18.0, 9.2, 7.0, 2.4, fill=PALE_TEAL, line=TEAL)
    add_text(s, "事前解析・熱真空試験で\n固定できる可能性", 18.45, 9.65, 6.1, 1.3, size=27, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "共有aの安定性がモデルの核心。符号は太陽面、絶対値はほぼ同じ。Adaptiveでaを毎回学習する必要性は低そう。")

    # 23
    s = new_slide(prs, "ケースDC b は太陽面＋発熱フラグで LOO RMSE 3.8 µrad", "Hierarchical ΔT model")
    img = ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/bcase_b_emp_vs_b_pred.png"
    add_picture_fit(s, img, 1.2, 3.0, 15.5, 8.6)
    add_table(s, ["Feature", "µrad"], [("b0_MX", "+15.7"), ("b0_MY", "+2.8"), ("b0_PX", "−12.0"), ("b0_PY", "−24.0"), ("c_PROP", "−22.9"), ("c_PCDU", "−10.2")], 17.4, 3.15, 7.7, 6.8, widths=[4.1, 3.6], font_size=23, accent_col=1)
    add_metric(s, "3.80 µrad", "LOO b RMSE", 17.5, 10.25, 7.5, TEAL)
    add_notes(s, "発熱効果はほぼ足し算。最大LOO残差は半電力case22で13.2 µrad。ON/OFFモデルの限界が明確。")

    # 24
    s = new_slide(prs, "標準ケースは数µradまで予測、半電力ケースでDCずれが残る", "Hierarchical ΔT model")
    img1 = ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/timeseries/case08_bcase_true_vs_pred.png"
    img2 = ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/timeseries/case22_bcase_true_vs_pred.png"
    add_picture_fit(s, img1, 1.25, 3.0, 11.9, 8.3)
    add_picture_fit(s, img2, 13.5, 3.0, 11.9, 8.3)
    add_label(s, "Case 08: standard / large swing", 2.8, 11.55, 8.8, TEAL)
    add_label(s, "Case 22: PROP 12.5 W", 15.6, 11.55, 7.7, ORANGE)
    add_text(s, "Level 1の傾きは維持されるが、ON/OFF Level 2では中間電力のbを表現できない", 3.2, 12.65, 20.4, 0.65, size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "良いケースと限界ケースを並べる。半電力では時変aは合うがbがずれるため、将来は発熱W比例へ拡張する。")

    # 25
    s = new_slide(prs, "生LOS数百µradに対し、LOO test残差は平均約5.5 µrad", "Hierarchical ΔT model")
    img = ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/bcase_raw_vs_model_rmse.png"
    add_picture_fit(s, img, 1.2, 3.0, 18.0, 8.7)
    add_metric(s, "615 µrad", "median raw RMS", 19.7, 3.5, 5.3, BLUE)
    add_metric(s, "5.5 µrad", "mean LOO test RMSE", 19.7, 6.0, 5.3, TEAL)
    add_metric(s, "1–2 orders", "thermal-bias reduction", 19.7, 8.5, 5.3, ORANGE)
    add_text(s, "raw RMSとprediction RMSEは異なる統計量：オーダー比較として表示", 3.6, 12.25, 19.5, 0.7, size=25, color=MID, align=PP_ALIGN.CENTER)
    add_notes(s, "モデル精度の主数字。raw RMSとRMSEを同じ意味で比較しないことを注記する。")

    # 26
    s = new_slide(prs, "新規性は式形ではなく、衛星バス相対LOSと粗捕捉への接続", "Positioning & limitations")
    add_box(s, 1.3, 3.0, 11.7, 8.3, fill=PALE_BLUE, line=BLUE)
    add_text(s, "JANUS / prior art", 1.8, 3.45, 10.7, 0.7, size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["Optical head内で LOS ≈ K·ΔT", "温度センサによる一次補正", "比例関係自体は既知"], 2.0, 4.7, 10.2, 3.4, size=27, bullet_color=BLUE)
    add_text(s, "式に切片を足すこと自体を\n新規性とは主張しない", 2.1, 8.7, 10.0, 1.3, size=27, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_box(s, 13.6, 3.0, 11.7, 8.3, fill=PALE_TEAL, line=TEAL)
    add_text(s, "This work", 14.1, 3.45, 10.7, 0.7, size=30, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["STT–LCTの衛星バス相対LOS", "太陽面−反対面ΔTの共有感度", "発熱DCを階層モデル化", "scan-center FFと捕捉時間へ接続"], 14.3, 4.7, 10.2, 4.5, size=27)
    add_text(s, "限界：MZ・非支配軸・別構造・飛行実証は未評価", 3.8, 12.1, 19.1, 0.8, size=28, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "新規性を慎重に置く。ΔT一次関係はJANUSにある。本研究の差分はバス相対LOS、階層DC、粗捕捉性能への接続。")

    # 27
    s = new_slide(prs, "予測LOSを scan centerから引き、残差をスキャン対象とする", "PAT evaluation")
    add_box(s, 1.5, 3.1, 23.6, 2.6, fill=PALE_TEAL, line=TEAL)
    add_text(s, "escan(t) = enonthermal(t) + θthermal,true(t) − θ̂thermal(t)", 2.0, 3.85, 22.6, 1.0, size=40, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_process(s, ["Thermal LOS\ntruth", "Model\nprediction", "Scan-center\nresidual", "Rectangular\nspiral", "Acquisition\ntime"], 1.5, 6.4, 23.6)
    add_table(s, ["Comparison arm", "Purpose"], [("no correction", "thermal baseline"), ("static bias", "DC-only baseline"), ("bcase", "proposed thermal model"), ("thermal truth", "ideal upper bound"), ("thermal+nonthermal, no corr.", "realism baseline"), ("bcase+nonthermal", "main result")], 4.1, 9.0, 18.6, 3.9, widths=[9.3, 9.3], font_size=20, accent_col=0)
    add_notes(s, "PATへ入れる式。非熱誤差は補正対象外で残る。比較armは熱モデル能力とシステム性能を分けて読む。")

    # 28
    s = new_slide(prs, "粗捕捉は矩形スパイラルを離散点で抽象化", "PAT simulation conditions")
    # spiral schematic
    add_box(s, 1.4, 3.0, 10.0, 8.5, fill=LIGHT, line=GRID)
    cx, cy = 6.4, 7.2
    pts = [(0,0),(1,0),(1,1),(-1,1),(-1,-1),(2,-1),(2,2),(-2,2),(-2,-2),(3,-2),(3,3)]
    scale = 0.75
    for i in range(len(pts)-1):
        add_arrow(s, cx+pts[i][0]*scale, cy-pts[i][1]*scale, cx+pts[i+1][0]*scale, cy-pts[i+1][1]*scale, color=BLUE, width=2)
    for px, py in pts:
        dot = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, I(cx+px*scale-0.09), I(cy-py*scale-0.09), I(0.18), I(0.18)); set_fill(dot, TEAL); dot.line.fill.background()
    target = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, I(cx+1.6), I(cy-1.8), I(0.55), I(0.55)); set_fill(target, PALE_ORANGE); set_line(target, ORANGE, 2)
    add_text(s, "target", cx+2.2, cy-1.85, 1.4, 0.4, size=18, color=ORANGE, bold=True)
    add_table(s, ["Parameter", "Value"], [("Range", "±1600 µrad"), ("Step", "40 µrad"), ("Detection radius", "25 µrad"), ("Dwell", "0.1 s / point"), ("Points", "6561"), ("Max nominal time", "656.1 s")], 12.2, 3.0, 12.9, 6.9, widths=[6.7, 6.2], font_size=23, accent_col=1)
    add_box(s, 12.2, 10.3, 12.9, 2.3, fill=PALE_RED, line=RED)
    add_text(s, "仮定：点間移動・settling・光強度確率なし\n各時刻を独立な捕捉機会として評価", 12.7, 10.75, 11.9, 1.3, size=25, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "モデルの単純化を丁寧に説明。平均捕捉時間は成功試行のみの条件付き平均で、成功率を必ず併記する。40 µrad格子と半径25 µradでは格子間に小さな未被覆領域もある。")

    # 29
    s = new_slide(prs, "非熱誤差は軌道・アライメント・姿勢・ドリフトを合成", "PAT simulation conditions")
    add_box(s, 1.5, 3.0, 23.6, 2.3, fill=PALE_BLUE, line=BLUE)
    add_text(s, "enonthermal = eorbit + ealignment + eattitude + edrift", 2.0, 3.65, 22.6, 0.9, size=39, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    rows = [("Orbit prediction", "Sentinel-1 TLE vs POEORB", "orbit-period family", "mean 280–700 µrad"), ("Alignment", "Gaussian constant / case", "DC", "50 µrad/axis, 1σ"), ("Attitude", "Gaussian / sample", "simplified broadband", "50 µrad/axis, 1σ"), ("Low-frequency drift", "sinusoid", "900 s", "30 µrad amplitude")]
    add_table(s, ["Component", "Model", "Timescale", "Size"], rows, 1.4, 6.0, 23.8, 5.2, widths=[4.7, 7.3, 5.3, 6.5], font_size=22, accent_col=0)
    add_text(s, "特定衛星の誤差バジェット再現ではなく、GNSS非搭載LEO小型衛星の共存シナリオ", 3.1, 12.05, 20.5, 0.8, size=28, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "軌道だけ実データ根拠が強い。姿勢・アライメント・ドリフトは代表オーダーの簡易モデル。")

    # 30
    s = new_slide(prs, "軌道誤差：利用可能な最新TLEをforward propagationし、PODと比較", "Orbit prediction error")
    add_process(s, ["Sentinel-1\nNORAD 39634", "Latest TLE\nepoch ≤ t", "SGP4 forward\nrpred(t)", "AUX_POEORB\nrtruth(t)", "δr(t)"], 1.3, 3.1, 24.0)
    img = ROOT / "results/orbit/sentinel1_tle_vs_pod/orbit_prediction_error_3orbits.png"
    add_picture_fit(s, img, 1.35, 5.4, 15.2, 6.8)
    add_box(s, 17.2, 5.4, 7.9, 6.8, fill=LIGHT, line=GRID)
    add_bullets(s, ["POEORBをtruth", "未来TLEのbackward伝搬は不使用", "60 s sampling", "相手機は片側truth", "熱解析時刻へcyclic resample"], 17.65, 5.95, 7.0, 4.6, size=24, bullet_color=BLUE)
    add_text(s, "2026 S1データとTD軌道時刻は位相一致していない", 17.7, 10.85, 6.9, 0.8, size=21, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "軌道誤差の作り方。予測として未来TLEを使わない。TDと実データ時刻は一致しないため、現状は代表外乱としてcyclicに載せている。")

    # 31
    s = new_slide(prs, "位置誤差をリンクLOS横断面へ投影し、STT/body x–yへ変換", "Orbit prediction error")
    add_box(s, 1.4, 3.0, 10.3, 8.4, fill=PALE_BLUE, line=BLUE)
    add_text(s, "Geometry", 1.9, 3.45, 9.3, 0.7, size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "δr⊥ = δr − (δr·ℓ̂)ℓ̂", 2.1, 5.0, 8.9, 0.9, size=31, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "θ ≈ δr⊥ / Rlink", 2.5, 6.5, 8.1, 0.9, size=33, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["orbit catalogの太陽面＋第二軸から姿勢を構成", "LOS横断誤差をSTT/body x–yへ射影", "熱LOSと同じ2成分としてPATへ加算"], 2.0, 8.0, 9.2, 2.8, size=23, bullet_color=BLUE)
    add_table(s, ["Sun case", "Partner / link", "Realism"], [("MY", "along-track ISL, 800 km", "ready"), ("PY", "anti-along-track ISL, 800 km", "ready"), ("PX", "nadir ground, ~695 km", "ready"), ("MX", "zenith proxy, 800 km", "unrealistic")], 12.4, 3.0, 12.8, 7.0, widths=[2.8, 6.8, 3.2], font_size=22, accent_col=2)
    add_box(s, 12.4, 10.4, 12.8, 2.0, fill=PALE_RED, line=RED)
    add_text(s, "MXは角度計算は可能だが現実的リンクではない\n→ 全ケース平均の解釈に注意", 12.9, 10.78, 11.8, 1.2, size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "以前の任意ECEF-Z基底ではなくSTT/body座標に接続した。MXだけLCTが天頂向きとなりリンクproxyが非現実的。")

    # 32
    s = new_slide(prs, "熱と軌道予測誤差はともに 10⁻⁴ Hz台の軌道周期族", "Orbit prediction error")
    img = ROOT / "results/orbit/sentinel1_tle_vs_pod/orbit_error_stt_LTAN06_800km_1213COLD_MY_SUN_3orbits.png"
    add_picture_fit(s, img, 1.25, 3.0, 15.7, 8.6)
    add_table(s, ["Case", "mean norm"], [("MY ISL", "≈285 µrad"), ("PY ISL", "≈285 µrad"), ("PX ground", "≈694 µrad"), ("MX proxy", "≈612 µrad")], 17.6, 3.2, 7.4, 5.4, widths=[3.8, 3.6], font_size=23, accent_col=1)
    add_box(s, 17.6, 9.0, 7.4, 2.6, fill=PALE_ORANGE, line=ORANGE)
    add_text(s, "Thermal: ~101 min\nOrbit norm: ~101 / 50 min", 18.0, 9.45, 6.6, 1.3, size=25, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "単純な『低周波＝熱』分離は難しい", 4.3, 12.15, 18.0, 0.8, size=31, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "五十里先生の指摘への回答。熱も軌道も軌道位相依存で帯域が重なる。従来の低周波成分を熱として取り込むAdaptive説明は修正が必要。")

    # 33
    s = new_slide(prs, "熱のみ：階層モデルは thermal-truth上界にほぼ到達", "PAT results")
    add_horizontal_bars(s, ["No correction", "Static bias", "Hierarchical bcase", "Thermal truth"], [124.6, 4.78, 0.116, 0.100], 1.6, 3.4, 15.7, 6.5, colors=[RED, ORANGE, TEAL, BLUE], suffix=" s", max_value=125)
    add_table(s, ["Model", "Success", "Mean tacq"], [("No correction", "95.9%", "124.6 s"), ("Static bias", "99.3%", "4.78 s"), ("bcase", "100%", "0.116 s"), ("truth", "100%", "0.100 s")], 18.0, 3.3, 7.1, 5.4, widths=[3.2, 2.0, 1.9], font_size=21, accent_col=2)
    add_box(s, 18.0, 9.1, 7.1, 2.7, fill=PALE_TEAL, line=TEAL)
    add_text(s, "熱残差平均\n8.8 µrad", 18.4, 9.55, 6.3, 1.4, size=31, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "熱モデル自体の能力を見る理想条件。システム主結果は次ページの非熱込み。", 3.0, 12.35, 20.7, 0.7, size=26, color=MID, align=PP_ALIGN.CENTER)
    add_notes(s, "Staticが非常に良いのは多くのケースでDC支配だから。ただし捕捉時刻や条件変化へは時変モデルが安定。平均時間は成功試行条件付き。")

    # 34
    s = new_slide(prs, "非熱込み：平均捕捉時間156.9→59.6 s、成功率94.0→97.1%", "PAT results")
    add_horizontal_bars(s, ["No thermal correction", "Hierarchical bcase FF"], [156.9, 59.6], 1.7, 3.5, 14.8, 4.2, colors=[RED, TEAL], suffix=" s", max_value=160)
    add_metric(s, "−62%", "mean acquisition time", 17.5, 3.4, 7.2, TEAL)
    add_metric(s, "+3.1 pt", "success rate", 17.5, 5.8, 7.2, BLUE)
    add_table(s, ["Metric", "No correction", "bcase FF"], [("Success", "94.0%", "97.1%"), ("Mean tacq", "156.9 s", "59.6 s"), ("Median tacq", "150.3 s", "36.2 s"), ("Mean thermal residual", "667 µrad", "8.8 µrad")], 3.1, 8.4, 20.5, 3.6, widths=[8.0, 6.25, 6.25], font_size=22, accent_col=2)
    add_text(s, "補正後は非熱、特に軌道予測誤差が性能を支配", 5.0, 12.55, 16.7, 0.7, size=30, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "主結果。21ケース単純平均。MY/PYでは大幅改善、PX/MXは軌道誤差が大きく改善率が小さい。MX proxyを含むため平均の解釈に注意。")

    # 35
    s = new_slide(prs, "Adaptiveは『低周波熱抽出』ではなく、主にモデルDC更新として考える", "Discussion")
    add_box(s, 1.3, 3.0, 11.8, 8.3, fill=PALE_TEAL, line=TEAL)
    add_text(s, "Proposed role", 1.8, 3.45, 10.8, 0.7, size=30, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "b̂k+1 = b̂k + Kb rk", 2.1, 4.85, 10.2, 0.9, size=36, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["a：事前解析／熱真空試験で固定", "b_case：捕捉後残差で更新", "同一運用モード・複数パスで平均", "軌道誤差やalignmentの吸収を制約"], 2.0, 6.3, 10.3, 4.1, size=25)
    add_box(s, 13.7, 3.0, 11.6, 8.3, fill=PALE_ORANGE, line=ORANGE)
    add_text(s, "Discussion points", 14.2, 3.45, 10.6, 0.7, size=30, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, ["a固定・b更新の分担は妥当か", "同帯域の熱／軌道誤差を何で識別するか", "軌道誤差を別状態として同時推定すべきか", "現在のscan modelで十分か", "ICSOまでにAdaptiveをどこまで実装するか"], 14.4, 4.7, 10.2, 5.7, size=24, bullet_color=ORANGE)
    add_text(s, "予測可能な熱成分はかなり除去できた。次は軌道上でDCをどう校正し、非熱誤差と共存させるか。", 2.1, 12.05, 22.5, 1.0, size=29, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "ここから30分の議論へ。周波数分離が難しいので、運用モード、温度センサ、軌道予測の事前情報を使った状態推定が必要。")

    # Appendix slides
    appendix_specs = [
        ("関連研究の位置づけ", [
            ("Badás et al. 2024", "FSOC opto-thermo-mechanical review"),
            ("Shi et al. 2023", "bus structure optimization → acquisition time"),
            ("Hu et al. 2022", "GEO thermal-motion Fourier compensation"),
            ("Li et al. 2025", "LEO geometry → NN LOS correction"),
            ("JANUS", "optical-head LOS ≈ KΔT"),
            ("Rüddenklau et al.", "body-pointing FF, not thermal deformation"),
        ]),
        ("全解析ケースの設計", []),
        ("熱光学特性", [("Alodine 1000", "α=0.150, ε=0.038, α/ε=3.947"), ("Black", "α=0.974, ε=0.920, α/ε=1.059"), ("0.5 synthetic", "α=0.500, ε=0.500")]),
        ("LOS定義の比較", []),
        ("Level-2係数と解釈", [("b0_MX", "+15.7 µrad"), ("b0_MY", "+2.8 µrad"), ("b0_PX", "−12.0 µrad"), ("b0_PY", "−24.0 µrad"), ("c_PROP", "−22.9 µrad"), ("c_PCDU", "−10.2 µrad")]),
        ("モデルの難しいケース", []),
        ("軌道予測誤差の仮定", [("Truth", "Sentinel-1 AUX_POEORB"), ("Prediction", "latest TLE + SGP4 forward"), ("Partner", "truth-known, one-sided error"), ("Range", "800 km ISL or nadir altitude"), ("Resampling", "cyclic to TD/Femap time")]),
        ("STT-frame射影結果", [("MY", "ISL along-track / 285 µrad"), ("PY", "ISL anti-along-track / 285 µrad"), ("PX", "ground nadir / 694 µrad"), ("MX", "zenith proxy / 612 µrad")]),
        ("PATスキャンの注意", [("Coverage", "40 µrad grid, 25 µrad detect radius"), ("Dynamics", "no slew / settling model"), ("Target", "constant during one scan"), ("Statistics", "mean tacq conditional on success"), ("Opportunities", "each thermal time sample independent")]),
        ("Adaptive候補", [("Scalar EWMA", "b only; simplest"), ("Kalman filter", "b + orbit/alignment states"), ("Mode-conditioned table", "sun face × power mode"), ("Batch pass update", "avoid single-acquisition contamination")]),
        ("主要参考文献", []),
    ]

    for idx, (title, rows) in enumerate(appendix_specs, 1):
        s = new_slide(prs, title, "Appendix", appendix=True)
        if title == "全解析ケースの設計":
            add_table(s, ["Cases", "Sun", "Power / condition"], [
                ("04, 13–15", "MY", "ALL / +PROP / +PCDU / STTLCT"),
                ("08, 16, 18, 19", "PY", "ALL / STTLCT / +PROP / +PCDU"),
                ("05, 06, 20, 21", "PX", "ALL / STTLCT / +PROP / +PCDU"),
                ("09, 17, 23, 24", "MX", "ALL / STTLCT / +PROP / +PCDU"),
                ("10–12", "MY", "HOT / Black / Alodine"),
                ("22", "MY", "PROP 12.5 W"),
                ("25", "MY", "LTAN18 / 693 km"),
            ], 1.5, 3.0, 23.5, 8.0, widths=[4.0, 3.0, 16.5], font_size=24, accent_col=2)
        elif title == "LOS定義の比較":
            add_picture_fit(s, ROOT / "results/femap_deformation/15_LTAN06_800km_1213COLD_MY_STTLCT_HEAT_MY_0p5/los_definition_comparison.png", 1.4, 3.0, 23.8, 9.3)
        elif title == "モデルの難しいケース":
            paths = [
                ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/timeseries/case11_bcase_true_vs_pred.png",
                ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/timeseries/case22_bcase_true_vs_pred.png",
                ROOT / "results/pat_acquisition/sunface_deltaT_bcase_los_model/timeseries/case25_bcase_true_vs_pred.png",
            ]
            for j, p in enumerate(paths):
                add_picture_fit(s, p, 1.0 + j * 8.4, 3.1, 7.8, 7.2)
            add_table(s, ["Case", "Issue"], [("11 Black", "oracle floor ≈13 µrad"), ("22 half power", "ON/OFF b mismatch"), ("25 all-sun", "thermal LOS nearly DC")], 3.0, 10.8, 20.6, 2.0, widths=[6.0, 14.6], font_size=21, accent_col=1)
        elif title == "主要参考文献":
            refs = [
                "[1] Riesing et al., On-orbit results of PAT for TBIRD, Proc. SPIE, 2023.",
                "[2] Shi et al., Thermal Deformation Stability Optimization..., Applied Sciences, 2023.",
                "[3] Badás et al., Opto-thermo-mechanical phenomena in satellite FSOC, Optical Engineering, 2024.",
                "[4] Hu et al., On-Board Thermal Motion Compensation..., IEEE GRSL, 2022.",
                "[5] Li et al., Correction Method for Thermal Deformation LOS Errors..., Remote Sensing, 2025.",
                "[6] Rüddenklau et al., Feed-Forward Compensation of Body-Pointing Uncertainties..., 2024.",
                "[7] JANUS optical-head thermo-elastic LOS model documentation / paper.",
            ]
            add_bullets(s, refs, 1.5, 3.1, 23.4, 8.8, size=25, bullet_color=PURPLE, gap=10)
        else:
            add_table(s, ["Item", "Detail"], rows, 2.0, 3.2, 22.6, min(8.5, 1.1 * (len(rows) + 1)), widths=[7.0, 15.6], font_size=24, header_color=PURPLE, accent_col=1)
        add_notes(s, f"Appendix: {title}. 質疑時の補足用。")

    return prs


def validate(prs: Presentation) -> None:
    problems = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                problems.append(f"slide {i}: negative position")
            if shape.left + shape.width > prs.slide_width + I(0.05):
                problems.append(f"slide {i}: shape exceeds width")
            if shape.top + shape.height > prs.slide_height + I(0.05):
                problems.append(f"slide {i}: shape exceeds height")
    if problems:
        raise RuntimeError("\n".join(problems[:20]))


def main() -> None:
    prs = build_deck()
    validate(prs)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
